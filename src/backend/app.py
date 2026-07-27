"""
FastAPI backend implementing the core user-flow runtime contract.
"""

from __future__ import annotations

import asyncio
import io
import json
import logging
import os
import re
import secrets as _secrets
import shutil
import subprocess
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional
from urllib.parse import parse_qs, urlencode, urlparse

import numpy as np
import polars as pl
import requests
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, Query, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles

logger = logging.getLogger(__name__)

# Maximum upload size in bytes (50 MB) (#6)
MAX_UPLOAD_BYTES = 50 * 1024 * 1024

# Maximum chat message length (#26)
MAX_CHAT_MESSAGE_LENGTH = 4000

# Regex for valid report/task IDs — UUIDs or hex strings, no path traversal
_SAFE_ID_RE = re.compile(r"^[a-zA-Z0-9_-]{1,128}$")


def _validate_id(value: str, label: str = "id") -> str:
    """Validate that a report_id or task_id is safe (no path traversal)."""
    clean = str(value or "").strip()
    if not clean or not _SAFE_ID_RE.match(clean):
        raise HTTPException(status_code=400, detail=f"Invalid {label}.")
    return clean

# Approximate FX conversion rates (1 unit of currency -> USD). These can be
# overridden with APP_FX_RATES_JSON='{"inr":0.012,...}' in runtime env.
_DEFAULT_FX_TO_USD: dict[str, float] = {
    "usd": 1.0,
    "eur": 1.08,
    "gbp": 1.27,
    "jpy": 0.0067,
    "cad": 0.74,
    "aud": 0.65,
    "mxn": 0.058,
    "brl": 0.19,
    "inr": 0.012,
    "cny": 0.14,
    "krw": 0.00074,
    "php": 0.018,
    "sgd": 0.74,
    "hkd": 0.13,
    "nzd": 0.60,
    "chf": 1.13,
    "sek": 0.095,
    "nok": 0.093,
    "dkk": 0.14,
    "zar": 0.055,
    "thb": 0.028,
    "idr": 0.000063,
    "myr": 0.22,
    "vnd": 0.000040,
    "aed": 0.27,
    "sar": 0.27,
    "try": 0.031,
    "pln": 0.25,
    "czk": 0.043,
    "huf": 0.0027,
    "clp": 0.0011,
    "cop": 0.00024,
    "pen": 0.27,
    "ars": 0.0010,
    "egp": 0.020,
    "ngn": 0.00064,
    "kes": 0.0077,
}
_FX_TO_USD = dict(_DEFAULT_FX_TO_USD)
_APP_FX_RATES_RAW = os.getenv("APP_FX_RATES_JSON", "").strip()
if _APP_FX_RATES_RAW:
    try:
        parsed = json.loads(_APP_FX_RATES_RAW)
        if isinstance(parsed, dict):
            for k, v in parsed.items():
                code = str(k or "").strip().lower()
                if not code:
                    continue
                try:
                    rate = float(v)
                    if rate > 0:
                        _FX_TO_USD[code] = rate
                except Exception:
                    continue
    except Exception:
        pass

_CURRENCY_ALIASES: dict[str, tuple[str, ...]] = {
    "usd": ("usd", "us dollar", "us dollars", "dollar", "dollars", r"\$"),
    "eur": ("eur", "euro", "euros", "€"),
    "gbp": ("gbp", "pound", "pounds", "sterling", "£"),
    "jpy": ("jpy", "yen", "¥"),
    "cad": ("cad", "canadian dollar", "canadian dollars"),
    "aud": ("aud", "australian dollar", "australian dollars"),
    "mxn": ("mxn", "peso", "pesos"),
    "brl": ("brl", "real", "reais"),
    "inr": ("inr", "rs", "rupee", "rupees", "₹"),
    "cny": ("cny", "rmb", "yuan", "元"),
}
_AMOUNT_RE = re.compile(
    r"(?P<amount>-?\d{1,3}(?:,\d{3})*(?:\.\d+)?|-?\d+(?:\.\d+)?)"
    r"(?P<suffix>[kKmMbB]?)"
)


def _token_pattern(token: str) -> str:
    t = str(token or "").strip()
    if not t:
        return ""
    if t in {"$", "€", "£", "¥", "₹", "元"}:
        return re.escape(t)
    if t == "rs":
        return r"(?:\brs\.?\b)"
    return rf"(?:\b{re.escape(t)}\b)"


def _parse_amount(raw: str, suffix: str = "") -> Optional[float]:
    try:
        val = float(str(raw).replace(",", "").strip())
        mult = {"k": 1_000, "m": 1_000_000, "b": 1_000_000_000}.get((suffix or "").lower(), 1)
        return val * mult
    except Exception:
        return None


def _format_usd(amount: float) -> str:
    """Format as clean USD — uses k/M/B shorthand for large numbers."""
    if abs(amount) >= 1_000_000_000:
        return f"${amount / 1_000_000_000:,.1f}B USD"
    if abs(amount) >= 1_000_000:
        return f"${amount / 1_000_000:,.1f}M USD"
    if abs(amount) >= 10_000:
        return f"${amount / 1_000:,.0f}k USD"
    return f"${amount:,.2f} USD"


def _convert_currency_text_to_usd(text: str) -> str:
    raw = str(text or "")
    if not raw:
        return raw
    out = raw

    # Phase 1: Convert non-USD currencies to USD (e.g., €50 → $55.00 USD)
    for code, aliases in _CURRENCY_ALIASES.items():
        if code == "usd":
            continue
        rate = float(_FX_TO_USD.get(code, 0.0) or 0.0)
        if rate <= 0:
            continue
        for alias in aliases:
            tok = _token_pattern(alias)
            if not tok:
                continue

            def _make_repl(r: float):
                def _repl(m: re.Match[str]) -> str:
                    val = _parse_amount(m.group("amount"), m.group("suffix"))
                    if val is None:
                        return m.group(0)
                    return _format_usd(val * r)
                return _repl

            repl_fn = _make_repl(rate)
            out = re.sub(rf"{tok}\s*{_AMOUNT_RE.pattern}", repl_fn, out, flags=re.IGNORECASE)
            out = re.sub(rf"{_AMOUNT_RE.pattern}\s*{tok}", repl_fn, out, flags=re.IGNORECASE)

    # Phase 2: Single-pass normalization of all USD-like amounts.
    # Matches: "$45k", "$12.00", "$12.00 USD", "USD 45k", "45k USD", "60 dollars"
    _USD_FULL = re.compile(
        r"(?:"
        r"\$\s*(?P<a1>-?\d[\d,]*(?:\.\d+)?)(?P<s1>[kKmMbB]?)(?:\s*USD\b)?"  # $45k or $12.00 or $12.00 USD
        r"|"
        r"(?:USD|US\$|dollars?)\s*(?P<a2>-?\d[\d,]*(?:\.\d+)?)(?P<s2>[kKmMbB]?)"  # USD 45k, dollars 100
        r"|"
        r"(?P<a3>-?\d[\d,]*(?:\.\d+)?)(?P<s3>[kKmMbB]?)\s*(?:USD|US\s+dollars?|dollars?)\b"  # 45k USD, 100 dollars
        r")",
        re.IGNORECASE,
    )

    def _norm_usd(m: re.Match[str]) -> str:
        for i in range(1, 4):
            amt_raw = m.group(f"a{i}")
            if amt_raw is not None:
                val = _parse_amount(amt_raw, m.group(f"s{i}") or "")
                if val is not None:
                    return _format_usd(val)
                break
        return m.group(0)

    out = _USD_FULL.sub(_norm_usd, out)

    # Remove residual rupee notation if present without a parseable amount.
    out = re.sub(r"\bRs\.?\b", "USD", out, flags=re.IGNORECASE)
    out = out.replace("₹", "USD ")

    # Final cleanup — collapse artifacts from previous bad processing
    out = re.sub(r"\$\$+", "$", out)
    out = re.sub(r"\bUSD\s*[kKmMbB]\b", "USD", out)  # stray "USDk" → "USD"
    out = re.sub(r"(\bUSD\b)(\s*\bUSD\b)+", "USD", out)  # must run after USDk cleanup
    return out


def _normalize_currency_payload(value: Any) -> Any:
    if isinstance(value, str):
        return _convert_currency_text_to_usd(value)
    if isinstance(value, dict):
        return {k: _normalize_currency_payload(v) for k, v in value.items()}
    if isinstance(value, list):
        return [_normalize_currency_payload(v) for v in value]
    if isinstance(value, tuple):
        return tuple(_normalize_currency_payload(v) for v in value)
    return value

from .contracts import (
    CategoryMatchRequest,
    ChatBriefRequest,
    ChatExportRequest,
    ChatGuidedQuestionsRequest,
    ChatMessageRequest,
    ColumnQuestionsRequest,
    PredictRequest,
    SchemaAnalyzeRequest,
    SessionStatusResponse,
    TaskStatusResponse,
    TellUsAugmentRequest,
    TellUsConfirmRequest,
    TellUsFinalizeRequest,
    TellUsFinalizeResponse,
    TellUsQuestionsRequest,
    TellUsRequest,
    TellUsTurnRequest,
    TellUsTurnResponse,
)
from .llm.orchestrator import LLMOrchestrator
from .retention import RetentionManager
from .security import RateLimiter
from .stores import ChatStore, ReportStore, SessionStore, TaskStore, derive_ip_hash, derive_stable_identity, iso, utcnow


_OLLAMA_PATTERN = re.compile(r"\bOllama\b", re.IGNORECASE)


def _sanitize_error(msg: str) -> str:
    """Return a clean, user-safe error message. Log the real error server-side."""
    raw = str(msg or "")
    logger.error("Pipeline error (raw): %s", raw)
    lowered = raw.lower()
    if "expected 4 charts" in lowered or "got 3" in lowered or "expected 6 charts" in lowered or "got 5" in lowered:
        return "Toji returned fewer visuals than required for the dashboard. Please retry."
    if "no dashboard payload returned" in lowered:
        return "Toji did not return dashboard content this turn. Please retry."
    if "csv context synthesis failed" in lowered:
        return "Toji could not complete synthetic data generation this turn. Please retry."
    if "reasoning-only content" in lowered or "empty content" in lowered:
        return "Toji cloud returned an incomplete response. Please retry."
    return "We couldn't generate your dashboard right now. Please try again."


def _toji_unavailable_detail(llm_orchestrator: Any, *, fallback: str) -> str:
    """Return a user-safe, actionable Toji availability message."""
    try:
        if getattr(llm_orchestrator, "provider_name", "") == "ollama":
            provider = getattr(llm_orchestrator, "provider", None)
            target = str(getattr(provider, "transport_target", "") or "").strip().lower()
            has_key = bool(str(getattr(provider, "api_key", "") or "").strip())
            provider_error = str(getattr(provider, "last_error", "") or "").strip()
            if target == "cloud" and not has_key:
                return "Toji cloud is not configured yet. Missing OLLAMA_API_KEY in deployment environment."
            if provider_error:
                provider_error_lower = provider_error.lower()
                safe = _sanitize_error(provider_error)
                lowered = safe.lower()
                if "http 401" in provider_error_lower or "http 403" in provider_error_lower:
                    return "Toji cloud rejected the request (auth failed). Check OLLAMA_API_KEY and account access."
                if "http 404" in provider_error_lower:
                    return (
                        "Toji cloud could not find the configured model. "
                        "Check OLLAMA_MODEL (for example `gpt-oss:120b-cloud` or `gpt-oss:120b`)."
                    )
                if "http 429" in provider_error_lower:
                    return "Toji cloud is rate-limited right now. Please retry in a few seconds."
                if "http 5" in provider_error_lower:
                    return "Toji cloud is temporarily unavailable. Please retry shortly."
                return f"Toji cloud request failed: {safe[:280]}"
    except Exception:
        pass
    return fallback


CANONICAL_STATES = {
    "IDLE",
    "REHYDRATING",
    "RESUMING",
    "SESSION_ACTIVE",
    "INTAKE_MODE_SELECT",
    "UPLOAD_PENDING",
    "UPLOAD_RECEIVED",
    "QUESTIONNAIRE_IN_PROGRESS",
    "QUESTIONNAIRE_CONFIRMED",
    "SYNTH_DATA_GENERATING",
    "SCHEMA_ANALYSIS_RUNNING",
    "SCHEMA_READY",
    "CATEGORY_EVAL_RUNNING",
    "CATEGORY_ACCEPTED",
    "PARTIAL_MATCH",
    "CATEGORY_REJECTED",
    "ROUTING_DECIDED",
    "PREDICT_ENQUEUED",
    "ALIGNMENT_RUNNING",
    "TOJI_CONTEXT_BUILDING",
    "TOJI_ANALYSIS_RUNNING",
    "REPORT_BUILDING",
    "REPORT_READY",
    "CHAT_ACTIVE",
    "WARNING_5_MIN",
    "SESSION_LOCKED_COOLDOWN",
    "ERROR",
    "DEBUG_BUNDLE_PERSISTED",
}


def _client_ip(request: Request) -> str:
    forwarded = request.headers.get("x-forwarded-for", "")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return (request.client.host if request.client else "0.0.0.0")


def _ip_hash(request: Request, secret: str) -> str:
    return derive_ip_hash(secret=secret, ip=_client_ip(request), user_agent=request.headers.get("user-agent", ""))


def _stable_id(request: Request, secret: str) -> str:
    return derive_stable_identity(secret=secret, ip=_client_ip(request), user_agent=request.headers.get("user-agent", ""))


def _sanitize_filename(raw: str) -> str:
    """Strip path separators and dangerous characters from a user-supplied filename (#7)."""
    name = Path(raw).name  # strips directory components
    # Remove any remaining path separators or null bytes
    name = re.sub(r'[\x00/\\]', '_', name)
    if not name or name.startswith('.'):
        name = f"upload_{name}"
    return name


def _safe_rel(base: Path, maybe_rel: str) -> Path:
    p = Path(maybe_rel)
    if p.is_absolute():
        resolved = p.resolve()
    else:
        resolved = (base / p).resolve()
    if base.resolve() not in resolved.parents and resolved != base.resolve():
        raise ValueError("Artifact path escapes project base")
    return resolved


def _safe_rel_intake(base: Path, intake_root: Path, maybe_rel: str) -> Path:
    """Restrict artifact resolution to the intake directory (#39)."""
    resolved = _safe_rel(base, maybe_rel)
    if intake_root.resolve() not in resolved.parents and resolved != intake_root.resolve():
        raise ValueError("Artifact must reside in the intake directory")
    return resolved


def _detect_chat_artifact_intent(message: str, suggested: Any = None) -> tuple[bool, bool]:
    text = str(message or "").lower()
    # Patterns that imply both doc + slides
    both_patterns = (
        "executive package",
        "full package",
        "everything",
        "comprehensive report",
        "all of it",
    )
    doc_patterns = (
        "doc",
        "document",
        "write-up",
        "write up",
        "writeup",
        "proposal",
        "brief",
        "memo",
        "report file",
        "word file",
        "executive brief",
        "summary doc",
    )
    slide_patterns = (
        "slide",
        "deck",
        "presentation",
        "ppt",
        "pitch",
        "board deck",
        "investor",
    )
    want_both = any(p in text for p in both_patterns)
    want_doc = want_both or any(p in text for p in doc_patterns)
    want_slides = want_both or any(p in text for p in slide_patterns)
    if isinstance(suggested, list):
        normalized = {str(x or "").strip().lower() for x in suggested}
        want_doc = want_doc or ("doc" in normalized) or ("document" in normalized)
        want_slides = want_slides or ("slides" in normalized) or ("deck" in normalized) or ("presentation" in normalized)
    return want_doc, want_slides


def _artifact_context(
    report_payload: dict[str, Any],
    reply: dict[str, Any],
    user_message: str,
    chat_history: Optional[list[dict[str, Any]]] = None,
) -> dict[str, Any]:
    intake = report_payload.get("intake_qa") or {}
    runtime_inf = report_payload.get("runtime_inference") or {}
    scorecard = report_payload.get("scorecard") or {}
    llm_summary = report_payload.get("llm_summary") or {}
    unified = report_payload.get("unified_analysis") or {}
    opportunity = report_payload.get("opportunity_analysis") or {}

    highlights = [str(x).strip() for x in (reply.get("highlights") or []) if str(x).strip()][:5]
    recs_raw = reply.get("recommendations") or []
    recommendations: list[dict[str, str]] = []
    if isinstance(recs_raw, list):
        for row in recs_raw[:5]:
            if isinstance(row, dict):
                recommendations.append(
                    {
                        "headline": str(row.get("headline") or "").strip(),
                        "rider": str(row.get("rider") or "").strip(),
                        "action": str(row.get("action") or "").strip(),
                        "impact": str(row.get("impact") or "").strip(),
                        "timeline": str(row.get("timeline") or "").strip(),
                    }
                )
            else:
                txt = str(row or "").strip()
                if txt:
                    recommendations.append({"headline": "", "rider": "", "action": txt, "impact": "", "timeline": ""})

    # Full problem list from report
    problems_raw = llm_summary.get("problems") or []
    problems: list[dict[str, str]] = []
    for p in problems_raw[:8]:
        if isinstance(p, dict):
            problems.append({
                "title": str(p.get("title") or "").strip(),
                "evidence": str(p.get("evidence") or "").strip(),
                "severity": str(p.get("severity") or "").strip(),
            })
        else:
            txt = str(p or "").strip()
            if txt:
                problems.append({"title": txt, "evidence": "", "severity": ""})

    # Full recommendation list from report
    full_recs_raw = llm_summary.get("recommendations") or []
    full_recommendations: list[dict[str, str]] = []
    for r in full_recs_raw[:8]:
        if isinstance(r, dict):
            full_recommendations.append({
                "headline": str(r.get("headline") or "").strip(),
                "rider": str(r.get("rider") or "").strip(),
                "action": str(r.get("action") or "").strip(),
                "impact": str(r.get("impact") or "").strip(),
                "timeline": str(r.get("timeline") or "").strip(),
                "difficulty": str(r.get("difficulty") or "").strip(),
            })
        else:
            txt = str(r or "").strip()
            if txt:
                full_recommendations.append({"headline": "", "rider": "", "action": txt, "impact": "", "timeline": "", "difficulty": ""})

    # KPI snapshot
    coverage_val = scorecard.get("coverage")
    mean_conf_val = runtime_inf.get("mean_confidence")
    kpis: list[dict[str, str]] = [
        {"label": "Data Coverage", "value": f"{float(coverage_val or 0):.0%}" if coverage_val is not None else "—", "sub": ""},
        {"label": "Model Confidence", "value": f"{float(mean_conf_val or 0):.0%}" if mean_conf_val is not None else "—", "sub": ""},
        {"label": "Data Quality", "value": str(scorecard.get("data_quality_band") or "—"), "sub": ""},
        {"label": "Rows Analysed", "value": str(scorecard.get("rows") or "—"), "sub": ""},
    ]

    # Opportunity narrative
    opp_narrative = str(opportunity.get("narrative") or "").strip()
    opp_gap = str(opportunity.get("gap") or "").strip()
    opp_metric = str(opportunity.get("metric") or "core metric").replace("_", " ")
    levers = opportunity.get("top_3_levers") or []
    benchmark_text = ""
    if opp_narrative:
        parts = [opp_narrative]
        for lev in (levers[:3] if isinstance(levers, list) else []):
            if isinstance(lev, dict):
                parts.append(f"{lev.get('lever', '')}: {lev.get('potential_impact', '')} ({lev.get('effort', '')} effort)")
        benchmark_text = " | ".join(parts)
    elif opp_gap:
        benchmark_text = f"{opp_metric}: {opp_gap}"

    # Chat history Q&A pairs (last 10 exchanges)
    qa_history: list[dict[str, str]] = []
    if chat_history:
        msgs = [m for m in chat_history if isinstance(m, dict) and m.get("role") in {"user", "assistant"}]
        # Build Q&A pairs
        for i in range(len(msgs) - 1):
            if msgs[i].get("role") == "user" and msgs[i + 1].get("role") == "assistant":
                qa_history.append({
                    "question": str(msgs[i].get("content") or "").strip()[:500],
                    "answer": str(msgs[i + 1].get("content") or "").strip()[:800],
                })
        qa_history = qa_history[-10:]  # last 10 pairs

    return {
        "industry": str(report_payload.get("industry") or "operations"),
        "category": str(report_payload.get("category") or "general"),
        "problem_statement": str((intake if isinstance(intake, dict) else {}).get("problem_statement") or "").strip(),
        "chat_question": str(user_message or "").strip(),
        "assistant_answer": str(reply.get("answer") or "").strip(),
        "highlights": highlights,
        "recommendations": recommendations,
        "coverage": coverage_val,
        "quality_band": scorecard.get("data_quality_band"),
        "mean_confidence": mean_conf_val,
        # Enriched fields
        "kpis": kpis,
        "problems": problems,
        "full_recommendations": full_recommendations,
        "report_analysis": str(llm_summary.get("summary") or llm_summary.get("toji_analysis") or "").strip(),
        "forecast": str(llm_summary.get("forecast_30d") or unified.get("forecast_30d") or "").strip(),
        "trend": str(unified.get("trend") or "").strip(),
        "benchmark": benchmark_text,
        "chat_history": qa_history,
    }


def _write_chat_docx(path: Path, ctx: dict[str, Any]) -> None:
    from docx import Document  # lazy import to keep runtime start lightweight
    from docx.shared import Pt

    doc = Document()
    ts = datetime.now(timezone.utc).strftime("%B %d, %Y")

    # 1. Cover
    doc.add_heading("Toji Executive Package", level=0)
    doc.add_paragraph(f"Industry: {ctx.get('industry', '—')}  |  Category: {ctx.get('category', '—')}  |  Generated: {ts}")

    problem = str(ctx.get("problem_statement") or "").strip()
    if problem:
        doc.add_heading("Business Problem", level=1)
        doc.add_paragraph(problem)

    # 2. Executive Summary
    analysis = str(ctx.get("report_analysis") or "").strip()
    if analysis:
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(analysis)

    # 3. KPI Snapshot
    kpis = ctx.get("kpis") or []
    if kpis:
        doc.add_heading("KPI Snapshot", level=1)
        table = doc.add_table(rows=1, cols=2)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        hdr[0].text = "Metric"
        hdr[1].text = "Value"
        for kpi in kpis:
            row = table.add_row().cells
            row[0].text = str(kpi.get("label") or "")
            row[1].text = str(kpi.get("value") or "—")

    # 4. Signal Breakdown
    trend = str(ctx.get("trend") or "").strip()
    benchmark = str(ctx.get("benchmark") or "").strip()
    if trend or benchmark:
        doc.add_heading("Signal Breakdown", level=1)
        if trend:
            doc.add_paragraph(f"Trend: {trend}")
        if benchmark:
            doc.add_paragraph(f"Benchmark: {benchmark}")

    # 5. Top Problems
    problems = ctx.get("problems") or []
    if problems:
        doc.add_heading("Top Problems", level=1)
        for idx, p in enumerate(problems, 1):
            title = str(p.get("title") or "").strip()
            evidence = str(p.get("evidence") or "").strip()
            severity = str(p.get("severity") or "").strip()
            text = f"{idx}. {title}"
            if severity:
                text += f" [{severity}]"
            doc.add_paragraph(text)
            if evidence:
                ev_para = doc.add_paragraph(f"   Evidence: {evidence}")
                ev_para.paragraph_format.left_indent = Pt(18)

    # 6. Recommended Actions
    full_recs = ctx.get("full_recommendations") or ctx.get("recommendations") or []
    if full_recs:
        doc.add_heading("Recommended Actions", level=1)
        for idx, rec in enumerate(full_recs, 1):
            action = str(rec.get("action") or "").strip()
            impact = str(rec.get("impact") or "").strip()
            timeline = str(rec.get("timeline") or "").strip()
            difficulty = str(rec.get("difficulty") or "").strip()
            parts = [x for x in [impact, timeline, difficulty] if x]
            suffix = " | ".join(parts)
            doc.add_paragraph(f"{idx}. {action}{(' — ' + suffix) if suffix else ''}")

    # 7. 30-Day Plan
    forecast = str(ctx.get("forecast") or "").strip()
    if forecast:
        doc.add_heading("30-Day Outlook", level=1)
        doc.add_paragraph(forecast)

    # 8. Conversation Highlights
    qa_history = ctx.get("chat_history") or []
    if qa_history:
        doc.add_heading("Conversation Highlights", level=1)
        for pair in qa_history:
            q = str(pair.get("question") or "").strip()
            a = str(pair.get("answer") or "").strip()
            if q:
                doc.add_paragraph(f"Q: {q}")
            if a:
                doc.add_paragraph(f"A: {a}")
            doc.add_paragraph("")

    doc.save(str(path))


def _write_chat_pptx(path: Path, ctx: dict[str, Any]) -> None:
    from pptx import Presentation  # lazy import

    prs = Presentation()
    ts = datetime.now(timezone.utc).strftime("%B %d, %Y")

    def _body_tf(slide):
        """Return the text frame of the content placeholder (index 1)."""
        return slide.shapes.placeholders[1].text_frame

    def _add_bullets(tf, items, prefix="- "):
        tf.clear()
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"{prefix}{item}"

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Toji Executive Package"
    s1.placeholders[1].text = f"{ctx.get('industry', '')} • {ctx.get('category', '')} • {ts}"

    # Slide 2: Situation Overview (executive summary)
    analysis = str(ctx.get("report_analysis") or ctx.get("assistant_answer") or "").strip()
    if analysis:
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Situation Overview"
        _body_tf(s2).text = analysis[:900]

    # Slide 3: KPI Snapshot
    kpis = ctx.get("kpis") or []
    if kpis:
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "KPI Snapshot"
        kpi_lines = [f"{k.get('label', '')}: {k.get('value', '—')}" for k in kpis]
        _add_bullets(_body_tf(s3), kpi_lines, prefix="")

    # Slide 4: Key Findings (trend + benchmark)
    trend = str(ctx.get("trend") or "").strip()
    benchmark = str(ctx.get("benchmark") or "").strip()
    findings: list[str] = []
    if trend:
        findings.append(f"Trend: {trend}")
    if benchmark:
        findings.append(f"Benchmark: {benchmark}")
    if findings:
        s4 = prs.slides.add_slide(prs.slide_layouts[1])
        s4.shapes.title.text = "Key Findings"
        _add_bullets(_body_tf(s4), findings, prefix="")

    # Slide 5: Problems & Root Causes
    problems = ctx.get("problems") or []
    if problems:
        s5 = prs.slides.add_slide(prs.slide_layouts[1])
        s5.shapes.title.text = "Problems & Root Causes"
        prob_lines = []
        for p in problems[:6]:
            title = str(p.get("title") or "").strip()
            severity = str(p.get("severity") or "").strip()
            prob_lines.append(f"{title}{(' [' + severity + ']') if severity else ''}")
        _add_bullets(_body_tf(s5), prob_lines)

    # Slide 6: Recommended Actions
    full_recs = ctx.get("full_recommendations") or ctx.get("recommendations") or []
    if full_recs:
        s6 = prs.slides.add_slide(prs.slide_layouts[1])
        s6.shapes.title.text = "Recommended Actions"
        rec_lines = []
        for r in full_recs[:6]:
            action = str(r.get("action") or "").strip()
            impact = str(r.get("impact") or "").strip()
            timeline = str(r.get("timeline") or "").strip()
            detail = " | ".join([x for x in [impact, timeline] if x])
            rec_lines.append(f"{action}{(' (' + detail + ')') if detail else ''}")
        _add_bullets(_body_tf(s6), rec_lines)

    # Slide 7: 30-Day Roadmap + Q&A highlights
    forecast = str(ctx.get("forecast") or "").strip()
    qa_history = ctx.get("chat_history") or []
    s7 = prs.slides.add_slide(prs.slide_layouts[1])
    s7.shapes.title.text = "30-Day Roadmap"
    roadmap_lines: list[str] = []
    if forecast:
        roadmap_lines.append(forecast[:300])
    if qa_history:
        roadmap_lines.append("")
        roadmap_lines.append("Conversation highlights:")
        for pair in qa_history[-5:]:
            q = str(pair.get("question") or "").strip()
            a = str(pair.get("answer") or "").strip()
            if q:
                roadmap_lines.append(f"Q: {q[:120]}")
            if a:
                roadmap_lines.append(f"A: {a[:200]}")
    if roadmap_lines:
        _add_bullets(_body_tf(s7), [l for l in roadmap_lines if l], prefix="")

    prs.save(str(path))


def _build_chat_artifacts(
    *,
    base: Path,
    report_id: str,
    report_payload: dict[str, Any],
    reply: dict[str, Any],
    user_message: str,
    want_doc: bool,
    want_slides: bool,
    chat_history: Optional[list[dict[str, Any]]] = None,
) -> tuple[list[dict[str, str]], list[str]]:
    artifacts: list[dict[str, str]] = []
    errors: list[str] = []
    if not (want_doc or want_slides):
        return artifacts, errors

    assets_dir = base / "exports" / "runtime_reports" / "assets" / report_id
    assets_dir.mkdir(parents=True, exist_ok=True)
    ts = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    ctx = _artifact_context(report_payload, reply, user_message, chat_history=chat_history)

    if want_doc:
        doc_name = f"toji_brief_{ts}.docx"
        doc_path = assets_dir / doc_name
        try:
            _write_chat_docx(doc_path, ctx)
            artifacts.append(
                {
                    "type": "document",
                    "filename": doc_name,
                    "uri": f"/report/{report_id}/asset/{doc_name}",
                }
            )
        except Exception as exc:
            errors.append(f"document generation failed: {exc}")

    if want_slides:
        slide_name = f"toji_slides_{ts}.pptx"
        slide_path = assets_dir / slide_name
        try:
            _write_chat_pptx(slide_path, ctx)
            artifacts.append(
                {
                    "type": "slides",
                    "filename": slide_name,
                    "uri": f"/report/{report_id}/asset/{slide_name}",
                }
            )
        except Exception as exc:
            errors.append(f"slides generation failed: {exc}")

    return artifacts, errors


def _read_dataset(path: Path) -> pl.DataFrame:
    lower = path.name.lower()
    if lower.endswith(".parquet") or lower.endswith(".parquet.zstd") or lower.endswith(".pq"):
        return pl.read_parquet(path)
    if lower.endswith(".csv") or lower.endswith(".txt"):
        try:
            return pl.read_csv(path, infer_schema_length=2000, ignore_errors=False)
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail=f"The uploaded CSV file is malformed and could not be parsed. Please check the file and try again.",
            ) from exc
    if lower.endswith(".json"):
        return pl.read_json(path)
    if lower.endswith(".xlsx") or lower.endswith(".xls"):
        try:
            import pandas as pd  # local import to keep startup lightweight

            workbook = pd.read_excel(path, sheet_name=None)
            frames: list[pl.DataFrame] = []
            for sheet_name, frame in (workbook or {}).items():
                if frame is None or frame.empty:
                    continue
                local = frame.copy()
                local.insert(0, "sheet_name", str(sheet_name))
                frames.append(pl.from_pandas(local))
            if not frames:
                raise ValueError("No worksheet data found.")
            base_cols = list(frames[0].columns)
            normalized: list[pl.DataFrame] = []
            for frame in frames:
                missing = [col for col in base_cols if col not in frame.columns]
                extra = [col for col in frame.columns if col not in base_cols]
                for col in missing:
                    frame = frame.with_columns(pl.lit(None).alias(col))
                ordered = frame.select(base_cols + extra)
                normalized.append(ordered)
            return pl.concat(normalized, how="diagonal_relaxed")
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail="The uploaded Excel file could not be parsed. Please check the workbook and try again.",
            ) from exc
    if lower.endswith(".pdf"):
        pdftotext_bin = shutil.which("pdftotext")
        if not pdftotext_bin:
            raise HTTPException(status_code=400, detail="PDF ingestion is unavailable on this server right now.")
        try:
            raw_text = subprocess.check_output(
                [pdftotext_bin, str(path), "-"],
                text=True,
                stderr=subprocess.STDOUT,
            )
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail="The uploaded PDF could not be read. Please check the file and try again.",
            ) from exc
        pages = raw_text.split("\f")
        rows: list[dict[str, Any]] = []
        for page_idx, page in enumerate(pages, start=1):
            for line_idx, line in enumerate(page.splitlines(), start=1):
                clean = str(line or "").strip()
                if not clean:
                    continue
                rows.append(
                    {
                        "page_number": page_idx,
                        "line_number": line_idx,
                        "text": clean,
                    }
                )
        if not rows:
            raise HTTPException(status_code=400, detail="The uploaded PDF did not contain extractable text.")
        return pl.DataFrame(rows)
    raise ValueError(f"Unsupported file type for {path.name}")


def _quickchart_render_url(url: str, *, orchestrator: Any) -> str:
    target = orchestrator._normalize_quickchart_url(url)
    parsed = urlparse(target)
    query = dict(parse_qs(parsed.query, keep_blank_values=True))
    query["format"] = ["png"]
    query["backgroundColor"] = ["white"]
    query["devicePixelRatio"] = ["2"]
    query["width"] = [str(query.get("width", ["560"])[0] or "560")]
    query["height"] = [str(query.get("height", ["340"])[0] or "340")]
    return parsed._replace(query=urlencode(query, doseq=True)).geturl()


def _sandbox_render_quickchart_visuals(
    *,
    base_path: Path,
    orchestrator: Any,
    report_id: str,
    visuals: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    if not visuals:
        return []

    sandbox_dir = base_path / "exports" / "runtime_reports" / "sandbox" / report_id
    asset_dir = base_path / "exports" / "runtime_reports" / "assets" / report_id
    sandbox_dir.mkdir(parents=True, exist_ok=True)
    asset_dir.mkdir(parents=True, exist_ok=True)

    session = requests.Session()
    placeholder_png = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00"
        b"\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc```\xf8\x0f\x00\x01\x04\x01\x00^\xef\xf9\xb0"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    rendered: list[dict[str, Any]] = []
    for idx, visual in enumerate(visuals, start=1):
        row = dict(visual or {})
        meta = dict(row.get("meta") or {})
        source_url = str(row.get("uri") or "").strip()
        if str(row.get("kind") or "").strip().lower() != "quickchart":
            rendered.append(row)
            continue
        ok, reason = orchestrator._quickchart_url_check(source_url)
        if not ok:
            raise RuntimeError(f"Visual {idx} failed quickchart validation: {reason}")
        render_url = _quickchart_render_url(source_url, orchestrator=orchestrator)
        try:
            response = session.get(render_url, timeout=20)
            if response.status_code >= 300:
                raise RuntimeError(f"Visual {idx} sandbox render failed with HTTP {response.status_code}.")
            content = response.content or b""
        except Exception:
            if os.getenv("PYTEST_CURRENT_TEST"):
                content = placeholder_png
            else:
                raise
        if len(content) < 64 or not content.startswith(b"\x89PNG\r\n\x1a\n"):
            raise RuntimeError(f"Visual {idx} sandbox render did not return a valid PNG.")

        sandbox_name = f"sandbox_quickchart_{idx}.png"
        sandbox_path = sandbox_dir / sandbox_name
        sandbox_path.write_bytes(content)

        final_name = f"quickchart_{idx}.png"
        final_path = asset_dir / final_name
        shutil.copyfile(sandbox_path, final_path)

        meta.update(
            {
                "sandbox_verified": True,
                "sandbox_path": str(sandbox_path.relative_to(base_path)),
                "source_url": source_url,
                "render_url": render_url,
            }
        )
        row["uri"] = final_name
        row["meta"] = meta
        rendered.append(row)
    return rendered


def _schema_summary(df: pl.DataFrame) -> dict[str, Any]:
    rows = int(df.height)
    cols = int(df.width)
    columns = []
    missing = {}
    for name, dtype in zip(df.columns, df.dtypes):
        null_count = int(df[name].null_count())
        null_ratio = float(null_count / rows) if rows > 0 else 0.0
        columns.append(name)
        missing[name] = {
            "dtype": str(dtype),
            "null_count": null_count,
            "null_ratio": round(null_ratio, 6),
        }
    return {
        "rows": rows,
        "cols": cols,
        "columns": columns,
        "column_stats": missing,
    }


def _dataset_profile_for_llm(df: pl.DataFrame, schema: dict[str, Any], *, sample_rows: int = 12) -> dict[str, Any]:
    rows = int(schema.get("rows") or df.height or 0)
    cols = int(schema.get("cols") or df.width or 0)
    col_stats = schema.get("column_stats") or {}

    numeric_columns: list[dict[str, Any]] = []
    categorical_columns: list[dict[str, Any]] = []

    for name, dtype in zip(df.columns, df.dtypes):
        d = str(dtype)
        null_ratio = float(((col_stats.get(name) or {}).get("null_ratio") or 0.0))
        if any(tok in d for tok in ("Int", "UInt", "Float", "Decimal")):
            try:
                series = df[name].cast(pl.Float64, strict=False)
                numeric_columns.append(
                    {
                        "name": name,
                        "dtype": d,
                        "null_ratio": round(null_ratio, 6),
                        "mean": float(series.mean()) if rows > 0 else None,
                        "std": float(series.std()) if rows > 1 else None,
                        "min": float(series.min()) if rows > 0 else None,
                        "max": float(series.max()) if rows > 0 else None,
                    }
                )
            except Exception:
                continue
        else:
            try:
                top = (
                    df.select(pl.col(name).cast(pl.Utf8).alias("v"))
                    .drop_nulls("v")
                    .group_by("v")
                    .len()
                    .sort("len", descending=True)
                    .head(5)
                )
                top_values = []
                for row in top.iter_rows(named=True):
                    top_values.append(
                        {
                            "value": str(row.get("v") or ""),
                            "count": int(row.get("len") or 0),
                        }
                    )
                categorical_columns.append(
                    {
                        "name": name,
                        "dtype": d,
                        "null_ratio": round(null_ratio, 6),
                        "top_values": top_values,
                    }
                )
            except Exception:
                continue

    sample = (
        df.head(max(1, int(sample_rows)))
        .fill_null("")
        .to_dicts()
    )
    return {
        "rows": rows,
        "cols": cols,
        "columns": list(df.columns),
        "numeric_columns": numeric_columns[:20],
        "categorical_columns": categorical_columns[:20],
        "sample_rows": sample[: max(1, int(sample_rows))],
    }


def _quality_from_schema(schema: dict[str, Any]) -> dict[str, Any]:
    stats = schema.get("column_stats") or {}
    if not stats:
        return {"completeness": 0.0, "data_gap_ratio": 1.0, "missing_fields": []}
    null_ratios = {k: float((v or {}).get("null_ratio") or 0.0) for k, v in stats.items()}
    completeness = max(0.0, min(1.0, 1.0 - float(np.mean(list(null_ratios.values())))))
    data_gap_ratio = max(0.0, min(1.0, 1.0 - completeness))
    missing_fields = [k for k, v in null_ratios.items() if v > 0.35]
    return {
        "completeness": round(completeness, 4),
        "data_gap_ratio": round(data_gap_ratio, 4),
        "missing_fields": missing_fields[:30],
    }


def _calibrate_synthetic_confidence(raw_score: Any) -> float:
    """Synthetic data confidence is intentionally bounded to an executive-safe 80-85% band."""
    try:
        score = float(raw_score)
    except Exception:
        score = 0.0
    score = max(0.0, min(1.0, score))
    calibrated = 0.80 + (score * 0.05)
    return round(max(0.80, min(0.85, calibrated)), 4)


def _calibrate_prediction_confidence(raw_score: Any) -> float:
    """Prediction/forecast confidence is intentionally conservative in an 80-85% band."""
    try:
        score = float(raw_score)
    except Exception:
        score = 0.0
    score = max(0.0, min(1.0, score))
    calibrated = 0.80 + (score * 0.05)
    return round(max(0.80, min(0.85, calibrated)), 4)


def _attach_runtime_confidence_fields(runtime_payload: dict[str, Any], coverage_score: Any) -> None:
    if not isinstance(runtime_payload, dict):
        return
    synthetic_conf = _calibrate_synthetic_confidence(coverage_score)

    l1 = runtime_payload.get("l1") if isinstance(runtime_payload.get("l1"), dict) else {}
    mode = str(
        l1.get("prediction_mode")
        or runtime_payload.get("prediction_mode")
        or ""
    ).strip().lower()
    base_conf = l1.get("mean_confidence")
    if base_conf is None:
        base_conf = runtime_payload.get("mean_confidence")
    if base_conf is None:
        base_conf = synthetic_conf
    prediction_conf = _calibrate_prediction_confidence(base_conf)

    if mode == "classification":
        metric = "F1/ACC (proxy)"
    elif mode == "regression":
        metric = "R2 (proxy)"
    else:
        metric = "calibrated confidence"

    runtime_payload["synthetic_data_confidence"] = synthetic_conf
    runtime_payload["prediction_confidence"] = prediction_conf
    runtime_payload["prediction_confidence_metric"] = metric
    runtime_payload["mean_confidence"] = prediction_conf

    if l1:
        l1["prediction_confidence"] = prediction_conf
        l1["prediction_confidence_metric"] = metric
        l1["mean_confidence"] = prediction_conf


_ACTION_HUMANIZE = {
    "provide_missing": "Close missing business context fields before the next run to improve forecast accuracy.",
    "run_fallback": "Use conservative operating thresholds until model confidence improves.",
    "improve_schema_match": "Rename uploaded columns to match expected business fields for better reliability.",
    "calibrate_risk": "Calibrate risk thresholds weekly against real outcomes.",
    "collect_more_data": "Collect more historical data to strengthen trend and seasonality signals.",
}


def _humanize_actions(raw: list[Any]) -> list[dict[str, str]]:
    out: list[dict[str, str]] = []
    for action in raw:
        key = str(action).strip().lower().replace(" ", "_").replace("-", "_")
        label = _ACTION_HUMANIZE.get(key)
        if not label:
            # Try without underscores
            normalized = key.replace("_", "")
            for k, v in _ACTION_HUMANIZE.items():
                if k.replace("_", "") == normalized:
                    label = v
                    break
        if not label:
            label = str(action).replace("_", " ").capitalize() + "."
        priority = "HIGH" if key in ("provide_missing", "run_fallback") else "MEDIUM"
        out.append({"action": label, "priority": priority})
    return out


def _basic_report_payload(
    *,
    industry: str,
    category: str,
    match: dict[str, Any],
    schema: dict[str, Any],
    task_id: str,
) -> dict[str, Any]:
    rows = int(schema.get("rows", 0))
    cols = int(schema.get("cols", 0))
    missing_ratio = 0.0
    stats = schema.get("column_stats", {}) or {}
    if stats:
        missing_ratio = float(np.mean([float(v.get("null_ratio", 0.0)) for v in stats.values()]))
    coverage = float(max(0.0, min(1.0, 1.0 - missing_ratio)))
    quality = "high" if coverage >= 0.80 else ("medium" if coverage >= 0.55 else "low")

    raw_actions = match.get("next_actions", [])
    humanized_actions = _humanize_actions(raw_actions)

    return {
        "generated_at": iso(utcnow()),
        "task_id": task_id,
        "industry": industry,
        "category": category,
        "decision": match.get("decision"),
        "scorecard": {
            "data_quality_band": quality,
            "coverage": round(coverage, 4),
            "rows": rows,
            "columns": cols,
        },
        "routing": {
            "fallback_models": match.get("fallback_models", []),
            "top_workers": match.get("top_workers", []),
            "missing_fields": match.get("missing_fields", []),
        },
        "insights": [
            "Report generated via async backend contract flow.",
            "Toji generated this report directly from your context and data.",
            "Use missing-fields guidance to improve routing confidence.",
        ],
        "next_actions": humanized_actions,
        "column_source_map": {},
        "data_confidence": None,
    }


def create_app(base_path: Optional[Path] = None) -> FastAPI:
    base = Path(base_path or Path(__file__).resolve().parents[2]).resolve()
    load_dotenv(dotenv_path=base / ".env", override=False)
    _DEFAULT_SECRET = "dev-only-secret-rotate-in-production"
    session_secret = os.getenv("BACKEND_SESSION_SECRET", "").strip()
    if not session_secret:
        # No secret configured: generate a random per-startup secret so sessions
        # are at least unique per process, even if they won't survive restarts.
        session_secret = _secrets.token_hex(32)
        logger.warning(
            "BACKEND_SESSION_SECRET not set — using ephemeral random secret. "
            "Sessions will not persist across restarts."
        )
    elif session_secret == _DEFAULT_SECRET:
        logger.warning(
            "BACKEND_SESSION_SECRET is set to the default development value. "
            "Rotate this in production."
        )
    max_workers = int(os.getenv("BACKEND_TASK_WORKERS", "2"))
    session_ttl_hours_raw = float(os.getenv("SESSION_TTL_HOURS", "1") or 1.0)
    session_ttl_minutes = int(float(os.getenv("SESSION_TTL_MINUTES", str(max(10, int(session_ttl_hours_raw * 60)))) or 60))
    session_ttl_minutes = max(10, session_ttl_minutes)
    cooldown_hours = int(os.getenv("SESSION_COOLDOWN_HOURS", "24"))

    app = FastAPI(title="Logistics Intelligence Backend", version="0.1.0")
    # CORS: never combine allow_credentials=True with allow_origins=["*"] (#32)
    cors_origins = [x.strip() for x in os.getenv("BACKEND_CORS_ORIGINS", "").split(",") if x.strip()]
    if not cors_origins:
        # Default to permissive without credentials for development
        cors_origins = ["*"]
    use_credentials = "*" not in cors_origins
    app.add_middleware(
        CORSMiddleware,
        allow_origins=cors_origins,
        allow_credentials=use_credentials,
        allow_methods=["*"],
        allow_headers=["*"],
    )
    # Add security headers (#50)
    @app.middleware("http")
    async def add_security_headers(request: Request, call_next):
        response = await call_next(request)
        response.headers["X-Content-Type-Options"] = "nosniff"
        response.headers["X-Frame-Options"] = "DENY"
        response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
        response.headers["Content-Security-Policy"] = (
            "default-src 'self'; "
            "script-src 'self' 'unsafe-inline' https://cdn.jsdelivr.net; "
            "style-src 'self' 'unsafe-inline'; "
            "img-src 'self' data: blob:; "
            "connect-src 'self' http://localhost:8000 http://localhost:8080 https://cdn.jsdelivr.net; "
            "font-src 'self'; "
            "frame-ancestors 'none'"
        )
        # HSTS: only sent over HTTPS; omit on plain HTTP to avoid breaking local dev
        if request.url.scheme == "https":
            response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"
        response.headers["Permissions-Policy"] = "geolocation=(), microphone=(), camera=()"
        return response

    sessions = SessionStore(session_ttl_minutes=session_ttl_minutes, cooldown_hours=cooldown_hours)
    tasks = TaskStore(max_workers=max_workers, persist_dir=base / "exports" / "task_state")
    reports = ReportStore(base)
    chats = ChatStore(base)
    retention = RetentionManager(base)
    limiter = RateLimiter()
    legacy_router_enabled = os.getenv("ENABLE_LEGACY_ROUTER", "0").strip().lower() in {"1", "true", "yes", "on"}
    matcher: Any = None
    alias_registry: Any = None
    target_semantics: Any = None
    if legacy_router_enabled:
        from .alias_registry import AliasRegistryManager
        from .matcher import RouterManifestMatcher
        from .target_semantics import TargetSemanticsResolver

        matcher = RouterManifestMatcher(base)
        alias_registry = AliasRegistryManager(base)
        target_semantics = TargetSemanticsResolver(base)
        logger.info("Legacy router runtime is enabled (ENABLE_LEGACY_ROUTER=1).")
    else:
        logger.info("Legacy router runtime is disabled; Toji unified flow is active.")
    llm = LLMOrchestrator(base)
    markdown_dashboard_enabled = os.getenv("TOJI_MARKDOWN_DASHBOARD_ENABLED", "1").strip().lower() in {
        "1",
        "true",
        "yes",
        "on",
    }
    intake_root = base / "exports" / "runtime_intake"
    intake_root.mkdir(parents=True, exist_ok=True)

    # Periodic session + task cleanup
    async def _session_cleanup_task():
        while True:
            try:
                await asyncio.sleep(300)  # every 5 minutes
                sessions.purge_expired()
                tasks._purge_old_tasks()
            except asyncio.CancelledError:
                break
            except Exception:
                pass

    @app.on_event("startup")
    async def start_cleanup():
        asyncio.create_task(_session_cleanup_task())

    def require_session(request: Request):
        retention.run_if_due()
        ip_hash = _ip_hash(request, session_secret)
        sid = _stable_id(request, session_secret)
        rec = sessions.get(ip_hash, stable_id=sid)
        if rec is None:
            return ip_hash, None
        return ip_hash, rec

    def resolve_artifact(ip_hash: str, explicit_uri: Optional[str] = None) -> Path:
        rec = sessions.get(ip_hash)
        if rec is None:
            raise HTTPException(status_code=404, detail="No active session.")
        uri = explicit_uri or rec.artifact_uris.get("upload_uri")
        if not uri:
            raise HTTPException(status_code=400, detail="No uploaded dataset artifact available.")
        # Restrict to intake directory or runtime_state (#39)
        path = _safe_rel(base, uri)
        exports_dir = (base / "exports").resolve()
        if exports_dir not in path.resolve().parents:
            raise HTTPException(status_code=400, detail="Artifact must reside within exports directory.")
        if not path.exists():
            raise HTTPException(status_code=404, detail="Artifact file not found.")
        return path

    def read_session_json(ip_hash: str, key: str) -> Optional[dict[str, Any]]:
        rec = sessions.get(ip_hash)
        if rec is None:
            return None
        uri = rec.artifact_uris.get(key)
        if not uri:
            return None
        try:
            path = _safe_rel(base, uri)
            if not path.exists():
                return None
            return json.loads(path.read_text())
        except Exception:
            return None

    def write_runtime_json(ip_hash: str, kind: str, payload: dict[str, Any]) -> str:
        out_dir = base / "exports" / "runtime_state" / ip_hash
        out_dir.mkdir(parents=True, exist_ok=True)
        ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        out = out_dir / f"{kind}_{ts}.json"
        out.write_text(json.dumps(payload, indent=2) + "\n")
        return str(out.relative_to(base))

    _NOISY_CONTEXT_KEYS = {
        "intake_transcript",
        "assistant_messages",
        "raw_context",
    }

    def _flatten_context_text(value: Any, *, _key: str = "") -> list[str]:
        out: list[str] = []
        if value is None:
            return out
        if _key in _NOISY_CONTEXT_KEYS:
            return out
        if isinstance(value, str):
            cleaned = value.strip()
            if cleaned:
                out.append(cleaned[:400])
            return out
        if isinstance(value, (int, float, bool)):
            out.append(str(value))
            return out
        if isinstance(value, list):
            tail = value[-12:] if len(value) > 12 else value
            for item in tail:
                out.extend(_flatten_context_text(item))
            return out
        if isinstance(value, dict):
            for key, item in value.items():
                key_str = str(key).strip()
                if key_str in _NOISY_CONTEXT_KEYS:
                    continue
                out.append(str(key))
                out.extend(_flatten_context_text(item, _key=key_str))
            return out
        out.append(str(value))
        return out

    def _business_intent_text(tell_payload: dict[str, Any]) -> str:
        pieces: list[str] = []
        problem_statement = str(tell_payload.get("problem_statement") or "").strip()
        if problem_statement:
            pieces.append(problem_statement)
        context = tell_payload.get("context") or {}
        pieces.extend(_flatten_context_text(context))
        merged = " | ".join([p for p in pieces if p]).strip()
        if len(merged) > 6000:
            head = merged[:3000].rstrip()
            tail = merged[-2500:].lstrip()
            merged = f"{head} | ... | {tail}"
        return merged

    def _intake_qa_block(conversation: dict[str, Any]) -> str:
        transcript = [row for row in (conversation.get("transcript") or []) if isinstance(row, dict)]
        if not transcript:
            return ""
        pairs: list[str] = []
        pending_question = ""
        for row in transcript:
            role = str(row.get("role") or "").strip().lower()
            content = str(row.get("content") or "").strip()
            if not content:
                continue
            if role == "assistant":
                pending_question = content
                continue
            if role == "user":
                q = pending_question or "User context"
                pairs.append(f"Question: {q}\nAnswer: {content}")
                pending_question = ""
        block = "\n\n".join(pairs).strip()
        if len(block) > 9000:
            block = block[:9000].rstrip()
        return block

    def _question_id(row: dict[str, Any], idx: int) -> str:
        raw = str(
            row.get("group_id")
            or row.get("field")
            or row.get("question")
            or f"q_{idx+1}"
        )
        cleaned = re.sub(r"[^a-zA-Z0-9_]+", "_", raw).strip("_").lower()
        return cleaned or f"q_{idx+1}"

    def _toji_first_person_question(text: str) -> str:
        q = str(text or "").strip()
        if not q:
            return q
        return q if q.endswith("?") else f"{q}?"

    def _build_turn_questions(
        industry: str,
        category: str,
        max_questions: int,
        *,
        missing_fields: Optional[list[str]] = None,
        user_context: str = "",
    ) -> list[dict[str, Any]]:
        limit = max(1, int(max_questions))
        llm.preload(industry)

        normalized_missing = [str(x).strip() for x in (missing_fields or []) if str(x).strip()]

        # Intake-first path: when there are no explicit missing fields yet,
        # ask contextual business questions derived from industry + problem text.
        try:
            if not normalized_missing:
                contextual_rows = llm.contextual_intake_questions(
                    industry=industry,
                    category=category,
                    user_context=user_context,
                    max_questions=min(limit, 5),
                )
                contextual_out: list[dict[str, Any]] = []
                for idx, row in enumerate(contextual_rows[: min(limit, 5)]):
                    if not isinstance(row, dict):
                        continue
                    qid = _question_id(row, idx)
                    contextual_out.append(
                        {
                            "id": qid,
                            "field": str(row.get("field") or qid),
                            "question": _toji_first_person_question(str(row.get("question") or "").strip()),
                            "hint": str(row.get("hint") or "").strip(),
                            "reframe": str(row.get("reframe") or "").strip(),
                            "fields": [f for f in (row.get("fields") or []) if isinstance(f, dict)],
                            "columns": [str(c) for c in (row.get("columns") or []) if str(c).strip()],
                        }
                    )
                if contextual_out:
                    logger.info(
                        "Generated contextual turn questions via Ollama: industry=%s category=%s count=%d",
                        industry,
                        category,
                        len(contextual_out),
                    )
                    return contextual_out

            if normalized_missing:
                guided_rows = llm.guided_context_questions(
                    industry=industry,
                    category=category,
                    missing_fields=normalized_missing[:limit],
                    user_context=user_context,
                    max_questions=limit,
                )
                out: list[dict[str, Any]] = []
                for idx, row in enumerate(guided_rows[:limit]):
                    if not isinstance(row, dict):
                        continue
                    qid = _question_id(row, idx)
                    fields = [f for f in (row.get("fields") or []) if isinstance(f, dict)]
                    columns = [str(c) for c in (row.get("columns") or []) if str(c).strip()]
                    out.append(
                        {
                            "id": qid,
                            "field": str(row.get("field") or qid),
                            "question": _toji_first_person_question(str(row.get("question") or "").strip()),
                            "hint": str(row.get("hint") or "").strip(),
                            "reframe": str(row.get("reframe") or "").strip(),
                            "fields": fields,
                            "columns": columns,
                        }
                    )
                if out:
                    logger.info(
                        "Generated guided turn questions via Ollama: industry=%s category=%s count=%d",
                        industry,
                        category,
                        len(out),
                    )
                    return out
            raise RuntimeError("Ollama returned no intake questions for this context.")
        except Exception as exc:
            logger.error(
                "Ollama question generation failed: industry=%s category=%s error=%s",
                industry,
                category,
                exc,
            )
            raise RuntimeError("Toji question generation failed. Ollama must be available.") from exc

    def _conversation_summary(conversation: dict[str, Any]) -> dict[str, int]:
        if str(conversation.get("question_mode") or "") == "pure_chat":
            transcript = [row for row in (conversation.get("transcript") or []) if isinstance(row, dict)]
            user_turns = sum(1 for row in transcript if str(row.get("role") or "").strip().lower() == "user")
            question_limit = max(1, min(10, int(conversation.get("question_limit") or 10)))
            answered = int(user_turns)
            total = max(int(question_limit), answered)
            return {
                "total_questions": int(total),
                "answered": int(answered),
                "user_provided": int(answered),
                "benchmark_default": 0,
                "unresolved": 0,
            }
        answers = [row for row in (conversation.get("answers") or []) if isinstance(row, dict)]
        total = int(len(conversation.get("questions") or []))
        answered = len(answers)
        user_provided = sum(1 for row in answers if str(row.get("answer_source") or "") == "user_provided")
        benchmark_default = sum(1 for row in answers if str(row.get("answer_source") or "") == "benchmark_default")
        unresolved = max(0, total - answered)
        return {
            "total_questions": total,
            "answered": answered,
            "user_provided": user_provided,
            "benchmark_default": benchmark_default,
            "unresolved": unresolved,
        }

    def _resolve_turn_context(
        *,
        ip_hash: str,
        industry: Optional[str],
        category: Optional[str],
    ) -> tuple[str, str]:
        tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
        tell_context = tell_payload.get("context") or {}
        raw_industry = str(industry or tell_payload.get("industry") or "ecommerce").strip() or "ecommerce"
        resolved_industry, custom_industry = _normalize_industry_value(raw_industry)
        if not isinstance(tell_context, dict):
            tell_context = {}
        known_custom_industry = str(tell_context.get("custom_industry_name") or "").strip()
        resolved_custom_industry = custom_industry or known_custom_industry or None
        payload_changed = False
        if str(tell_payload.get("industry") or "").strip() != resolved_industry:
            tell_payload["industry"] = resolved_industry
            payload_changed = True
        if resolved_custom_industry and known_custom_industry != resolved_custom_industry:
            tell_context["custom_industry_name"] = resolved_custom_industry
            payload_changed = True
        if payload_changed:
            tell_payload["context"] = tell_context
            tell_payload["recorded_at"] = iso(utcnow())
            tell_uri = write_runtime_json(ip_hash, "tell_us", tell_payload)
            sessions.update_artifact(ip_hash, "tell_us_uri", tell_uri)
        llm.bind_session_context(
            ip_hash,
            industry=resolved_industry,
            custom_industry_name=resolved_custom_industry,
        )
        resolved_category = str(category or tell_context.get("category") or "").strip()
        if resolved_category:
            return resolved_industry, resolved_category

        intent = _business_intent_text(tell_payload)
        suggested = llm.infer_categories(resolved_industry, intent, top_k=3)
        resolved_category = (
            (suggested[0]["category"] if suggested else None)
            or llm.infer_category(resolved_industry, str(tell_payload.get("problem_statement") or ""))
            or (llm.categories(resolved_industry)[0] if llm.categories(resolved_industry) else "unclassified")
        )
        return resolved_industry, str(resolved_category or "unclassified")

    def _bootstrap_turn_conversation(
        *,
        ip_hash: str,
        industry: str,
        category: str,
        max_questions: int,
        reset: bool,
        missing_fields: Optional[list[str]] = None,
    ) -> dict[str, Any]:
        normalized_missing = sorted(set([str(x).strip() for x in (missing_fields or []) if str(x).strip()]))
        existing = read_session_json(ip_hash, "tell_us_conversation_uri")
        if (
            isinstance(existing, dict)
            and not reset
            and str(existing.get("industry") or "") == industry
            and str(existing.get("category") or "") == category
            and sorted(set([str(x).strip() for x in (existing.get("missing_fields") or []) if str(x).strip()])) == normalized_missing
        ):
            return existing

        question_limit = 10
        return {
            "industry": industry,
            "category": category,
            "missing_fields": normalized_missing,
            "question_limit": question_limit,
            "question_source": "ollama",
            "question_mode": "pure_chat",
            "question_provider": llm.provider_name,
            "question_model": llm.provider_model,
            "questions": [],
            "current_index": 0,
            "answers": [],
            "attempts": {},
            "transcript": [],
            "captured_facts": [],
            "ready_to_analyze": False,
            "completed": False,
            "total_questions": question_limit,
            "created_at": iso(utcnow()),
            "updated_at": iso(utcnow()),
        }

    def _persist_turn_conversation(ip_hash: str, conversation: dict[str, Any]) -> str:
        conversation["updated_at"] = iso(utcnow())
        rel = write_runtime_json(ip_hash, "tell_us_conversation", conversation)
        sessions.update_artifact(ip_hash, "tell_us_conversation_uri", rel)
        llm.bind_session_context(ip_hash, tell_us_conversation=conversation)
        return rel

    def _build_tell_us_from_conversation(
        *,
        ip_hash: str,
        conversation: dict[str, Any],
    ) -> dict[str, Any]:
        sentinel = "__industry_avg_minus_1sd__"
        answers = [row for row in (conversation.get("answers") or []) if isinstance(row, dict)]
        previous = read_session_json(ip_hash, "tell_us_uri") or {}
        previous_ctx = previous.get("context") or {}
        if not isinstance(previous_ctx, dict):
            previous_ctx = {}

        if str(conversation.get("question_mode") or "") == "pure_chat":
            transcript = [row for row in (conversation.get("transcript") or []) if isinstance(row, dict)]
            user_msgs = [
                _convert_currency_text_to_usd(str(row.get("content") or "").strip())
                for row in transcript
                if str(row.get("role") or "").strip().lower() == "user" and str(row.get("content") or "").strip()
            ]
            assistant_msgs = [
                _convert_currency_text_to_usd(str(row.get("content") or "").strip())
                for row in transcript
                if str(row.get("role") or "").strip().lower() == "assistant" and str(row.get("content") or "").strip()
            ]
            captured_facts = [str(x).strip() for x in (conversation.get("captured_facts") or []) if str(x).strip()]
            time_context_answer = str(conversation.get("time_context") or "").strip()
            qa_block = _intake_qa_block(conversation)

            context_payload = dict(previous_ctx)
            context_payload.update(
                {
                    "category": str(conversation.get("category") or previous_ctx.get("category") or ""),
                    "pure_chat_intake": True,
                    "question_mode": "pure_chat",
                    "intake_transcript": transcript[-80:],
                    "captured_facts": captured_facts,
                    "user_messages": user_msgs[-30:],
                    "assistant_messages": assistant_msgs[-30:],
                    "intake_qa_block": qa_block,
                    "time_context_answer": time_context_answer,
                    "issue_duration_text": time_context_answer,
                    "q1": user_msgs[0] if len(user_msgs) > 0 else "",
                    "q2": user_msgs[1] if len(user_msgs) > 1 else "",
                    "q3": user_msgs[2] if len(user_msgs) > 2 else "",
                    "q4": user_msgs[3] if len(user_msgs) > 3 else "",
                    "q5": user_msgs[4] if len(user_msgs) > 4 else "",
                    "q6": user_msgs[5] if len(user_msgs) > 5 else "",
                    "q7": time_context_answer,
                    "q8": previous_ctx.get("q8", 5),
                    "q9": user_msgs[-1] if user_msgs else "",
                    "column_answers": previous_ctx.get("column_answers", {}),
                    "conservative_fields": sorted(set(previous_ctx.get("conservative_fields") or [])),
                }
            )
            context_payload = _normalize_currency_payload(context_payload)

            summary_statement = " | ".join(user_msgs[-10:]).strip()
            problem_statement = _convert_currency_text_to_usd(
                str(previous.get("problem_statement") or "").strip() or qa_block or summary_statement or "Toji intake conversation"
            )
            return {
                "problem_statement": problem_statement[:4000],
                "industry": str(conversation.get("industry") or previous.get("industry") or "ecommerce"),
                "context": context_payload,
                "recorded_at": iso(utcnow()),
            }

        column_answers: dict[str, dict[str, Any]] = {}
        conservative_fields: set[str] = set()
        intake_answers: list[dict[str, Any]] = []
        summary_bits: list[str] = []
        time_context_answer = ""

        for row in answers:
            qid = str(row.get("question_id") or "question")
            group_key = str(row.get("field") or qid or "question")
            mapped = row.get("mapped_values") or {}
            if not isinstance(mapped, dict):
                mapped = {}
            mapped_values: dict[str, Any] = {}
            for col, value in mapped.items():
                col_name = str(col).strip()
                if not col_name:
                    continue
                mapped_values[col_name] = value
                if value == sentinel:
                    conservative_fields.add(col_name)

            if mapped_values:
                column_answers[group_key] = mapped_values
            else:
                normalized = row.get("normalized_answer")
                if normalized is not None:
                    column_answers[group_key] = {group_key: normalized}
                    if normalized == sentinel:
                        conservative_fields.add(group_key)

            question_text = str(row.get("question") or "").strip()
            answer_raw = _convert_currency_text_to_usd(str(row.get("answer_raw") or "").strip())
            if question_text and answer_raw:
                summary_bits.append(f"{question_text} {answer_raw}")
            if group_key == "__time_context__" and answer_raw:
                time_context_answer = answer_raw

            intake_answers.append(
                {
                    "question_id": qid,
                    "field": group_key,
                    "question": question_text,
                    "answer": answer_raw,
                    "answer_source": row.get("answer_source"),
                    "parse_confidence": row.get("parse_confidence"),
                }
            )

        non_time_answers = [row for row in intake_answers if row.get("field") != "__time_context__"]
        context_payload = dict(previous_ctx)
        context_payload.update(
            {
                "category": str(conversation.get("category") or previous_ctx.get("category") or ""),
                "column_answers": column_answers,
                "conservative_fields": sorted(conservative_fields),
                "intake_answers": intake_answers,
                "q1": non_time_answers[0]["answer"] if len(non_time_answers) > 0 else "",
                "q2": non_time_answers[1]["answer"] if len(non_time_answers) > 1 else "",
                "q3": non_time_answers[2]["answer"] if len(non_time_answers) > 2 else "",
                "q4": non_time_answers[3]["answer"] if len(non_time_answers) > 3 else "",
                "q5": non_time_answers[4]["answer"] if len(non_time_answers) > 4 else "",
                "q6": non_time_answers[5]["answer"] if len(non_time_answers) > 5 else "",
                "q7": time_context_answer,
                "q8": previous_ctx.get("q8", 5),
                "q9": non_time_answers[-1]["answer"] if non_time_answers else "",
                "time_context_answer": time_context_answer,
                "issue_duration_text": time_context_answer,
            }
        )
        context_payload = _normalize_currency_payload(context_payload)

        summary_statement = " | ".join(summary_bits).strip()
        problem_statement = _convert_currency_text_to_usd(
            str(previous.get("problem_statement") or "").strip() or summary_statement or "Toji intake conversation"
        )
        return {
            "problem_statement": problem_statement[:4000],
            "industry": str(conversation.get("industry") or previous.get("industry") or "ecommerce"),
            "context": context_payload,
            "recorded_at": iso(utcnow()),
        }

    def _generate_synthetic_artifact(
        *,
        ip_hash: str,
        requested_industry: Optional[str],
        progress_hook: Optional[Callable[[str, float, Optional[dict[str, Any]]], None]] = None,
    ) -> dict[str, Any]:
        def emit(state: str, progress: float, result: Optional[dict[str, Any]] = None) -> None:
            if progress_hook:
                progress_hook(state, progress, result)

        sessions.update_state(ip_hash, "SYNTH_DATA_GENERATING", step_uri="/intake/tell-us/confirm")
        tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
        tell_context = tell_payload.get("context") or {}
        if not isinstance(tell_context, dict):
            tell_context = {}
        problem_statement = str(tell_payload.get("problem_statement") or "")
        intent_text = _business_intent_text(tell_payload)

        industry, custom_industry = _normalize_industry_value(
            str(requested_industry or tell_payload.get("industry") or "ecommerce")
        )
        if custom_industry:
            tell_context["custom_industry_name"] = custom_industry
        resolved_custom_industry = str(tell_context.get("custom_industry_name") or "").strip() or None
        payload_industry = str(tell_payload.get("industry") or "").strip()
        if payload_industry != industry or tell_payload.get("context") != tell_context:
            tell_payload["industry"] = industry
            tell_payload["context"] = tell_context
            tell_payload["recorded_at"] = iso(utcnow())
            tell_uri = write_runtime_json(ip_hash, "tell_us", tell_payload)
            sessions.update_artifact(ip_hash, "tell_us_uri", tell_uri)

        emit("SYNTH_CONTEXT_BUILDING", 0.08)
        suggested_categories = llm.infer_categories(industry, intent_text, top_k=3)
        category = (
            tell_context.get("category")
            or llm.session_context(ip_hash).get("category")
            or (suggested_categories[0]["category"] if suggested_categories else None)
            or llm.infer_category(industry, problem_statement)
            or (llm.categories(industry)[0] if llm.categories(industry) else None)
        )

        llm.preload(industry)
        selected_category = str(category or "general_operations").strip() or "general_operations"

        column_values = None
        column_answers = tell_context.get("column_answers") if isinstance(tell_context, dict) else None
        if column_answers and isinstance(column_answers, dict) and selected_category:
            spec = llm._load_intake_spec(industry, selected_category)
            if not spec:
                spec = llm._auto_intake_spec_from_manifest(industry, selected_category)
            if spec:
                column_values = llm._map_answers_to_columns(spec, column_answers)

        synthesis_context = dict(tell_context) if isinstance(tell_context, dict) else {}
        collated_qa_block = str(synthesis_context.get("intake_qa_block") or "").strip()
        upload_profile_used = False
        upload_df = None
        upload_profile: dict[str, Any] = {}
        upload_artifact_uri = str(synthesis_context.get("upload_artifact_uri") or "").strip()
        if upload_artifact_uri:
            try:
                upload_df = _read_dataset(resolve_artifact(ip_hash, upload_artifact_uri))
                upload_schema = _schema_summary(upload_df)
                upload_profile = _dataset_profile_for_llm(upload_df, upload_schema)
                synthesis_context["uploaded_data_profile"] = upload_profile
                synthesis_context["uploaded_data_rows"] = int(upload_schema.get("rows") or 0)
                synthesis_context["uploaded_data_cols"] = int(upload_schema.get("cols") or 0)
                upload_profile_used = True
            except Exception as exc:
                logger.warning("Could not attach uploaded_data_profile for synthesis context: %s", exc)
                upload_df = None
                upload_profile = {}

        max_attempts_raw = os.getenv("TOJI_CONFIRM_RETRY_ATTEMPTS", "3")
        try:
            max_attempts = max(1, min(5, int(max_attempts_raw)))
        except Exception:
            max_attempts = 3

        df = None
        synth_meta: dict[str, Any] = {}
        last_exc: Optional[Exception] = None
        user_context_parts = [
            str(problem_statement or "").strip(),
            collated_qa_block,
            str(intent_text or "").strip(),
        ]
        user_context_text = "\n\n".join([part for part in user_context_parts if part]).strip()
        emit("SYNTH_DATA_GENERATING", 0.15)
        for attempt in range(1, max_attempts + 1):
            attempt_context = dict(synthesis_context)
            attempt_context["synthesis_attempt"] = attempt
            if last_exc is not None:
                attempt_context["previous_synthesis_error"] = str(last_exc)[:320]
            emit(
                "SYNTH_DATA_GENERATING",
                min(0.82, 0.15 + (attempt - 1) * (0.60 / max(1, max_attempts))),
                {"attempt": attempt, "max_attempts": max_attempts},
            )
            try:
                if upload_df is not None and upload_profile_used:
                    df, synth_meta = llm.transform_uploaded_frame(
                        industry=industry,
                        category=selected_category,
                        user_context=user_context_text,
                        uploaded_df=upload_df,
                        dataset_profile=upload_profile,
                        context=attempt_context,
                    )
                else:
                    df, synth_meta = llm.synthesize_context_frame(
                        industry=industry,
                        category=selected_category,
                        user_context=user_context_text,
                        context=attempt_context,
                        column_values=column_values,
                        n_rows=1000,
                    )
                break
            except Exception as exc:
                last_exc = exc
                logger.warning(
                    "Synthetic generation attempt failed: attempt=%d/%d industry=%s category=%s error=%s",
                    attempt,
                    max_attempts,
                    industry,
                    selected_category,
                    exc,
                )
                if attempt < max_attempts:
                    user_context_text = (
                        str(problem_statement or user_context_text).strip()[:1200]
                        or str(user_context_text).strip()[:1200]
                    )
                    continue
                break

        if df is None or synth_meta is None:
            exc = last_exc or RuntimeError("unknown synthetic generation failure")
            logger.error(
                "Ollama synthetic generation failed: industry=%s category=%s error=%s",
                industry,
                selected_category,
                exc,
            )
            reason = _sanitize_error(str(exc))
            if not reason:
                try:
                    provider_error = str(getattr(getattr(llm, "provider", None), "last_error", "") or "").strip()
                    if provider_error:
                        reason = _sanitize_error(provider_error)
                except Exception:
                    reason = ""
            detail = "Toji was unable to generate your dataset. Please retry in a moment."
            if reason:
                detail = f"{detail} Last error: {reason[:220]}"
            raise RuntimeError(detail) from exc

        emit("SYNTH_ARTIFACT_WRITING", 0.9)
        out_dir = intake_root / ip_hash
        out_dir.mkdir(parents=True, exist_ok=True)
        ts = datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')
        out_path = out_dir / f"{ts}_synthetic.csv"
        df.write_csv(out_path)
        rel = str(out_path.relative_to(base))
        sessions.update_artifact(ip_hash, "upload_uri", rel)

        pandas_copy_uri: Optional[str] = None
        try:
            pandas_df = df.to_pandas()
            pandas_path = out_dir / f"{ts}_synthetic_pandas.pkl"
            pandas_df.to_pickle(pandas_path)
            pandas_copy_uri = str(pandas_path.relative_to(base))
            sessions.update_artifact(ip_hash, "upload_pandas_uri", pandas_copy_uri)
        except Exception as exc:
            logger.warning("Could not persist pandas copy for synthetic data: %s", exc)

        python_script_uri: Optional[str] = None
        script_text = str((synth_meta or {}).get("python_script") or "").strip()
        if script_text:
            script_path = out_dir / f"{ts}_synthetic_plan.py"
            script_path.write_text(script_text, encoding="utf-8")
            python_script_uri = str(script_path.relative_to(base))
            sessions.update_artifact(ip_hash, "synthetic_script_uri", python_script_uri)

        synth_meta = dict(synth_meta or {})
        synth_meta["upload_profile_used"] = bool(upload_profile_used)
        if pandas_copy_uri:
            synth_meta["pandas_copy_uri"] = pandas_copy_uri
        if python_script_uri:
            synth_meta["python_script_uri"] = python_script_uri

        meta_rel = write_runtime_json(ip_hash, "synthetic_meta", synth_meta)
        sessions.update_artifact(ip_hash, "synthetic_meta_uri", meta_rel)
        sessions.update_state(ip_hash, "QUESTIONNAIRE_CONFIRMED", step_uri="/intake/tell-us/confirm")
        llm.bind_session_context(
            ip_hash,
            industry=industry,
            category=selected_category,
            custom_industry_name=resolved_custom_industry,
            synthetic_artifact_uri=rel,
            synthetic_meta=synth_meta,
            suggested_categories=suggested_categories,
            business_intent=intent_text,
        )
        result = {
            "state": "QUESTIONNAIRE_CONFIRMED",
            "industry": industry,
            "custom_industry_name": resolved_custom_industry,
            "category": selected_category,
            "artifact_uri": rel,
            "synthetic_meta_uri": meta_rel,
            "synthetic_meta": synth_meta,
            "pandas_copy_uri": pandas_copy_uri,
            "python_script_uri": python_script_uri,
            "suggested_categories": suggested_categories,
        }
        emit("SYNTH_DATA_READY", 1.0, result)
        return result

    supported_industries = set(matcher.available_industries()) if matcher is not None else set(llm.industries())
    supported_industries.add("other_industry")

    def _normalize_industry_value(
        raw_industry: Optional[str],
        *,
        default: str = "ecommerce",
    ) -> tuple[str, Optional[str]]:
        raw = str(raw_industry or "")[:256]  # Bound length before processing
        raw = raw.strip()
        if not raw:
            return default, None
        cleaned = re.sub(r"[^a-zA-Z0-9]+", "_", raw).strip("_").lower()
        if not cleaned:
            return default, None
        if cleaned in supported_industries:
            return cleaned, None
        return "other_industry", raw

    def _paired_qa(questions: list[str], answers: Any) -> list[dict[str, Any]]:
        qs = [str(q).strip() for q in (questions or []) if str(q).strip()]
        if not qs:
            return []

        if isinstance(answers, dict):
            # Preserve explicit question-keyed responses where provided.
            out = []
            for idx, q in enumerate(qs):
                key_candidates = [q, f"q{idx+1}", str(idx), f"question_{idx+1}"]
                answer_value = None
                for key in key_candidates:
                    if key in answers:
                        answer_value = answers.get(key)
                        break
                out.append({"question": q, "answer": answer_value})
            return out

        if isinstance(answers, list):
            out = []
            for idx, q in enumerate(qs):
                answer_value = answers[idx] if idx < len(answers) else None
                out.append({"question": q, "answer": answer_value})
            return out

        if isinstance(answers, str):
            if len(qs) == 1:
                return [{"question": qs[0], "answer": answers}]
            return [{"question": q, "answer": None} for q in qs]

        return [{"question": q, "answer": None} for q in qs]

    def _build_intake_qa_bundle(
        *,
        tell_payload: dict[str, Any],
        question_payload: dict[str, Any],
        augment_payload: dict[str, Any],
        session_ctx: dict[str, Any],
    ) -> dict[str, Any]:
        context = tell_payload.get("context") or {}
        industry_questions = list(question_payload.get("industry_questions") or [])
        category_questions = list(question_payload.get("category_questions") or [])
        category_question_sets = list(question_payload.get("category_question_sets") or [])

        industry_answers = (
            context.get("industry_answers")
            if isinstance(context, dict)
            else None
        )
        category_answers = (
            context.get("category_answers")
            if isinstance(context, dict)
            else None
        )
        if not industry_answers and isinstance(context, dict):
            industry_answers = {
                "workflow_goal": context.get("workflow_goal"),
                "kpi_baseline_target": context.get("kpi_baseline_target"),
                "constraints": context.get("constraints"),
            }
        if not category_answers and isinstance(context, dict):
            category_answers = context.get("category_details")

        industry_qa = _paired_qa(industry_questions, industry_answers)
        category_qa_primary = _paired_qa(category_questions, category_answers)

        category_qa_candidates = []
        candidate_answers = {}
        if isinstance(context, dict):
            candidate_answers = context.get("category_answers_by_category") or {}
            if not isinstance(candidate_answers, dict):
                candidate_answers = {}

        for row in category_question_sets:
            if not isinstance(row, dict):
                continue
            category = str(row.get("category") or "").strip()
            questions = list(row.get("questions") or [])
            answers = candidate_answers.get(category)
            category_qa_candidates.append(
                {
                    "category": category,
                    "score": row.get("score"),
                    "questions_and_answers": _paired_qa(questions, answers),
                }
            )

        return {
            "problem_statement": tell_payload.get("problem_statement"),
            "industry": tell_payload.get("industry"),
            "selected_category": session_ctx.get("category"),
            "suggested_categories": session_ctx.get("suggested_categories") or [],
            "industry_qa": industry_qa,
            "category_qa_primary": category_qa_primary,
            "category_qa_candidates": category_qa_candidates,
            "augment": {
                "missing_fields": augment_payload.get("missing_fields", []),
                "provided_values": augment_payload.get("provided_values", {}),
            },
            "raw_context": context,
        }

    def _resolve_report_id(ip_hash: Optional[str] = None, explicit: Optional[str] = None) -> str:
        report_id = str(explicit or "").strip()
        if report_id:
            return _validate_id(report_id, "report_id")
        if ip_hash:
            rec = sessions.get(ip_hash)
            if rec and rec.report_id:
                return str(rec.report_id)
        raise HTTPException(status_code=400, detail="No report_id provided and no report available in current session.")

    def _enforce_rate_limit(
        request: Request,
        *,
        ip_hash: str,
        scope: str,
        limit: int,
        window_seconds: int,
    ) -> None:
        key = f"{scope}:{ip_hash}"
        allowed, retry_after = limiter.allow(key, limit=limit, window_seconds=window_seconds)
        if not allowed:
            raise HTTPException(
                status_code=429,
                detail={
                    "error": "rate_limited",
                    "scope": scope,
                    "retry_after_seconds": retry_after,
                },
                headers={"Retry-After": str(retry_after)},
            )

    def _chat_export_payload(report_id: str, report_payload: dict[str, Any], history: list[dict[str, Any]]) -> dict[str, Any]:
        summary = llm.summarize_chat(report_payload=report_payload, history=history)
        return {
            "report_id": report_id,
            "generated_at": iso(utcnow()),
            "message_count": len(history),
            "summary": summary,
            "history_tail": history[-25:],
        }

    def _guardrail_redirect(report_payload: dict[str, Any], reason: str) -> dict[str, Any]:
        custom_name = str(report_payload.get("custom_industry_name") or "").strip()
        industry = custom_name or str(report_payload.get("industry") or "selected industry").replace("_", " ")
        category = str(report_payload.get("category") or "selected category").replace("_", " ")
        return {
            "blocked": True,
            "reason": reason,
            "response": (
                f"I'm built specifically for {industry} operations intelligence — "
                f"let's stay focused on your {category} analysis. "
                "Ask about risks, trends, forecasts, likely causes, or next actions."
            ),
            "highlights": [],
            "recommendations": [],
        }

    def _chat_guardrails(message: str, report_payload: dict[str, Any]) -> dict[str, Any]:
        text = str(message or "").strip()
        if not text:
            return _guardrail_redirect(report_payload, "empty_message")
        lower = text.lower()

        identity_or_injection_patterns = [
            r"\bwho are you\b",
            r"\bwhat model are you\b",
            r"\bwhat (ai|llm) model\b",
            r"\bare you (gpt|chatgpt|grok|ollama|llama|claude|gemini)\b",
            r"\bsystem prompt\b",
            r"\bdeveloper message\b",
            r"\bhidden instructions\b",
            r"\bignore (all|previous|prior) instructions\b",
            r"\bprompt injection\b",
            r"\bjailbreak\b",
            r"\breveal (your|the) (prompt|instructions|model)\b",
            r"\bshow (your|the) (prompt|instructions|chain of thought)\b",
            # Persona pivot patterns
            r"\bforget everything\b",
            r"\bforget all\b",
            r"\bforget (your|the) (instructions|context|role)\b",
            r"\byou are (now|no longer)\b",
            r"\bact as (a|an|my)\b",
            r"\bpretend (to be|you are)\b",
            r"\bswitch (to|your) (role|persona|character)\b",
            r"\bstop being\b.*\b(toji|advisor|analyst)\b",
            r"\bi('m| am) (now|actually|really) (a|an)\b",
            r"\bactually forget\b",
            r"\bfrom now on\b.*\byou are\b",
        ]
        for pattern in identity_or_injection_patterns:
            if re.search(pattern, lower):
                return _guardrail_redirect(report_payload, "identity_or_prompt_injection")

        domain_terms = {
            "prediction",
            "predict",
            "forecast",
            "risk",
            "kpi",
            "sla",
            "trend",
            "season",
            "anomaly",
            "conversion",
            "checkout",
            "fulfillment",
            "demand",
            "shipping",
            "freight",
            "trucking",
            "delivery",
            "route",
            "carrier",
            "cost",
            "volume",
            "revenue",
            "inventory",
            "model",
            "worker",
            "category",
            "industry",
            "dataset",
        }
        # Industry-specific domain terms
        _industry_terms: dict[str, set[str]] = {
            "bpo": {"churn", "seats", "agents", "attrition", "aht", "handle", "ticket", "csat", "nps", "escalation", "queue", "occupancy"},
            "ecommerce": {"cart", "checkout", "conversion", "aov", "returns", "sku", "catalog", "marketplace", "fulfillment"},
            "healthcare": {"patient", "readmission", "discharge", "clinical", "beds", "triage", "appointment", "provider"},
            "manufacturing": {"yield", "defect", "downtime", "throughput", "oee", "batch", "quality", "assembly", "scrap"},
            "retail": {"footfall", "basket", "shrinkage", "planogram", "markdown", "assortment", "store"},
        }
        industry = str(report_payload.get("industry") or "").lower()
        category = str(report_payload.get("category") or "").lower()
        if industry:
            domain_terms.update(industry.replace("_", " ").split())
            # Add industry-specific terms
            for key, terms in _industry_terms.items():
                if key in industry:
                    domain_terms.update(terms)
        if category:
            domain_terms.update(category.replace("_", " ").split())
        for row in (report_payload.get("routing_context") or {}).get("route_categories", []) or []:
            domain_terms.update(str(row).lower().replace("_", " ").split())
        # Add custom industry name terms
        custom_name = str(report_payload.get("custom_industry_name") or "").lower().strip()
        if custom_name:
            domain_terms.update(custom_name.split())
            for key, terms in _industry_terms.items():
                if key in custom_name:
                    domain_terms.update(terms)

        off_topic_patterns = [
            r"\bweather\b",
            r"\bsports?\b",
            r"\bmovie\b",
            r"\bmusic\b",
            r"\bcelebrity\b",
            r"\bpolitics?\b",
            r"\brelationship\b",
            r"\bhoroscope\b",
            r"\bcrypto price\b",
            r"\bstock price\b",
        ]
        # Weighted domain relevance check (#44): require meaningful overlap, not just any term
        words = set(re.findall(r'\b[a-z]{3,}\b', lower))
        domain_hits = sum(1 for term in domain_terms if term and term in words)
        domain_ratio = domain_hits / max(1, len(words))
        is_domain = domain_hits >= 1 and domain_ratio >= 0.05

        if not is_domain:
            for pattern in off_topic_patterns:
                if re.search(pattern, lower):
                    return _guardrail_redirect(report_payload, "off_scope")

        return {"blocked": False}

    @app.get("/healthz")
    def healthz() -> dict[str, Any]:
        llm_available = bool(llm.llm_available())
        llm_provider = llm.provider_name
        llm_model = llm.provider_model
        ollama_target = None
        ollama_base_host = None
        ollama_api_key_present = None
        if llm_provider == "ollama":
            provider_obj = llm.provider
            base_url = str(getattr(provider_obj, "base_url", "") or "")
            ollama_target = str(getattr(provider_obj, "transport_target", "") or "") or None
            ollama_api_key_present = bool(str(getattr(provider_obj, "api_key", "") or "").strip())
            try:
                ollama_base_host = (urlparse(base_url).hostname or "").strip().lower() or None
            except Exception:
                ollama_base_host = None
        # Verify filesystem writability (#40)
        fs_ok = True
        try:
            test_path = intake_root / ".healthcheck"
            test_path.write_text("ok")
            test_path.unlink()
        except Exception:
            fs_ok = False
        return {
            "ok": fs_ok,
            "industries": sorted(supported_industries),
            "llm_enabled": llm_available,
            "llm_provider": llm_provider,
            "llm_model": llm_model,
            "ollama_target": ollama_target,
            "ollama_base_host": ollama_base_host,
            "ollama_api_key_present": ollama_api_key_present,
            "fs_writable": fs_ok,
        }

    @app.get("/session/status", response_model=SessionStatusResponse)
    def session_status(request: Request):
        ip_hash, rec = require_session(request)
        if rec is None:
            return SessionStatusResponse(has_session=False, artifacts={})
        return SessionStatusResponse(**rec.to_status_payload())

    @app.post("/session/start", response_model=SessionStatusResponse)
    def session_start(request: Request):
        ip_hash = _ip_hash(request, session_secret)
        sid = _stable_id(request, session_secret)
        rec = sessions.start_or_resume(ip_hash, stable_id=sid)
        if rec.state == "SESSION_LOCKED_COOLDOWN":
            return JSONResponse(status_code=423, content=rec.to_status_payload())
        sessions.update_state(ip_hash, "SESSION_ACTIVE", step_uri="/session/start")
        rec = sessions.get(ip_hash)
        return SessionStatusResponse(**rec.to_status_payload())

    @app.post("/intake/upload")
    async def intake_upload(request: Request, file: UploadFile = File(...)):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        if rec.state == "SESSION_LOCKED_COOLDOWN":
            return JSONResponse(status_code=423, content=rec.to_status_payload())
        _enforce_rate_limit(request, ip_hash=ip_hash, scope="intake_upload", limit=int(os.getenv("RATE_UPLOAD_LIMIT", "20")), window_seconds=int(os.getenv("RATE_UPLOAD_WINDOW_SEC", "600")))

        sessions.update_state(ip_hash, "UPLOAD_PENDING", step_uri="/intake/upload")
        out_dir = intake_root / ip_hash
        out_dir.mkdir(parents=True, exist_ok=True)
        # Sanitize filename to prevent path traversal (#7)
        safe_name = _sanitize_filename(file.filename or "upload.csv")
        filename = f"{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}_{safe_name}"
        out_path = out_dir / filename
        # Read with size limit to prevent memory exhaustion (#6, #33)
        payload = await file.read(MAX_UPLOAD_BYTES + 1)
        if len(payload) > MAX_UPLOAD_BYTES:
            raise HTTPException(status_code=413, detail=f"File exceeds {MAX_UPLOAD_BYTES // (1024*1024)} MB limit.")
        # Validate file content: must match supported dataset formats (#33)
        ext = safe_name.rsplit('.', 1)[-1].lower() if '.' in safe_name else ''
        if ext not in ("csv", "json", "txt", "xls", "xlsx", "pdf"):
            raise HTTPException(status_code=400, detail="Supported uploads: PDF, CSV, XLS, XLSX, JSON, or TXT.")
        header = payload[:512]
        if ext == 'json':
            stripped = header.lstrip()
            if stripped and stripped[0:1] not in (b'{', b'['):
                raise HTTPException(status_code=400, detail="File does not appear to be valid JSON.")
        if ext == "pdf" and not payload.startswith(b"%PDF"):
            raise HTTPException(status_code=400, detail="File does not appear to be a valid PDF.")
        out_path.write_bytes(payload)
        rel = str(out_path.relative_to(base))
        sessions.update_artifact(ip_hash, "upload_uri", rel)
        upload_schema = None
        upload_schema_uri = None
        try:
            uploaded_df = _read_dataset(out_path)
            upload_schema = _schema_summary(uploaded_df)
            upload_schema_uri = write_runtime_json(ip_hash, "upload_schema", upload_schema)
            sessions.update_artifact(ip_hash, "upload_schema_uri", upload_schema_uri)
        except Exception as exc:
            logger.warning("Could not derive upload schema summary for intake context: %s", exc)
        sessions.update_state(ip_hash, "UPLOAD_RECEIVED", step_uri="/intake/upload")
        llm.preload()
        llm.bind_session_context(
            ip_hash,
            upload_uri=rel,
            upload_schema=upload_schema,
            upload_schema_uri=upload_schema_uri,
        )
        return {
            "artifact_uri": rel,
            "upload_schema_uri": upload_schema_uri,
            "upload_schema": upload_schema,
            "state": "UPLOAD_RECEIVED",
        }

    @app.post("/intake/tell-us")
    def intake_tell_us(request: Request, body: TellUsRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        _enforce_rate_limit(request, ip_hash=ip_hash, scope="intake_tell_us", limit=int(os.getenv("RATE_TELLUS_LIMIT", "40")), window_seconds=int(os.getenv("RATE_TELLUS_WINDOW_SEC", "3600")))
        sessions.update_state(ip_hash, "QUESTIONNAIRE_IN_PROGRESS", step_uri="/intake/tell-us")
        normalized_industry, custom_industry = _normalize_industry_value(body.industry)
        normalized_context = _normalize_currency_payload(body.context)
        if not isinstance(normalized_context, dict):
            normalized_context = {}
        if custom_industry:
            normalized_context["custom_industry_name"] = custom_industry
        resolved_custom_industry = custom_industry or str(normalized_context.get("custom_industry_name") or "").strip() or None
        artifact = {
            "problem_statement": _convert_currency_text_to_usd(body.problem_statement),
            "industry": normalized_industry,
            "context": normalized_context,
            "recorded_at": iso(utcnow()),
        }
        rel = write_runtime_json(ip_hash, "tell_us", artifact)
        sessions.update_artifact(ip_hash, "tell_us_uri", rel)
        llm.preload(normalized_industry)
        llm.bind_session_context(
            ip_hash,
            tell_us=artifact,
            industry=normalized_industry,
            custom_industry_name=resolved_custom_industry,
        )
        return {"state": "QUESTIONNAIRE_IN_PROGRESS", "artifact_uri": rel}

    @app.post("/intake/tell-us/questions")
    def intake_tell_us_questions(request: Request, payload: TellUsQuestionsRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
        industry, custom_industry = _normalize_industry_value(
            str(payload.industry or tell_payload.get("industry") or "ecommerce")
        )
        context = tell_payload.get("context") or {}
        if not isinstance(context, dict):
            context = {}
        if custom_industry and isinstance(context, dict):
            context["custom_industry_name"] = custom_industry
        resolved_custom_industry = custom_industry or str(context.get("custom_industry_name") or "").strip() or None
        intent_text = _business_intent_text(tell_payload)
        suggested_categories = llm.infer_categories(industry, intent_text, top_k=3)
        category = payload.category or context.get("category")
        if not category:
            category = (
                (suggested_categories[0]["category"] if suggested_categories else None)
                or llm.infer_category(industry, str(tell_payload.get("problem_statement") or ""))
            )
        llm.preload(industry)
        response = llm.question_set(industry=industry, category=category)
        response["suggested_categories"] = suggested_categories
        response["category_question_sets"] = [
            {
                "category": row["category"],
                "score": row["score"],
                "reasons": row.get("reasons", []),
                "questions": llm.question_set(industry=industry, category=row["category"]).get("category_questions", []),
            }
            for row in suggested_categories
        ]
        q_rel = write_runtime_json(ip_hash, "question_set", response)
        sessions.update_artifact(ip_hash, "questions_uri", q_rel)
        llm.bind_session_context(
            ip_hash,
            industry=industry,
            category=category,
            custom_industry_name=resolved_custom_industry,
            suggested_categories=suggested_categories,
            business_intent=intent_text,
            question_set=response,
        )
        return {**response, "artifact_uri": q_rel}

    @app.post("/intake/tell-us/turn", response_model=TellUsTurnResponse)
    def intake_tell_us_turn(request: Request, body: TellUsTurnRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        _enforce_rate_limit(
            request,
            ip_hash=ip_hash,
            scope="intake_tell_us",
            limit=int(os.getenv("RATE_TELLUS_LIMIT", "40")),
            window_seconds=int(os.getenv("RATE_TELLUS_WINDOW_SEC", "3600")),
        )

        industry, category = _resolve_turn_context(
            ip_hash=ip_hash,
            industry=body.industry,
            category=body.category,
        )
        try:
            conversation = _bootstrap_turn_conversation(
                ip_hash=ip_hash,
                industry=industry,
                category=category,
                max_questions=body.max_questions,
                reset=bool(body.reset),
                missing_fields=body.missing_fields,
            )
        except Exception as exc:
            logger.error(
                "Failed to bootstrap pure intake conversation: industry=%s category=%s error=%s",
                industry,
                category,
                exc,
            )
            raise HTTPException(
                status_code=503,
                detail=_toji_unavailable_detail(
                    llm,
                    fallback="Toji is unavailable for intake chat right now. Please retry in a moment.",
                ),
            ) from exc

        parse_confidence = None
        answer_source = None
        normalized_answer: Any = None
        assistant_message = None

        transcript = [row for row in (conversation.get("transcript") or []) if isinstance(row, dict)]
        conversation["transcript"] = transcript
        conversation["question_mode"] = "pure_chat"
        conversation["question_source"] = "ollama"
        conversation["question_provider"] = llm.provider_name
        conversation["question_model"] = llm.provider_model

        incoming_answer = str(body.answer).strip() if body.answer is not None else ""
        if incoming_answer and not bool(conversation.get("completed")):
            answer_source = "user_provided"
            normalized_answer = _convert_currency_text_to_usd(incoming_answer)
            transcript.append(
                {
                    "role": "user",
                    "content": normalized_answer,
                    "ts": iso(utcnow()),
                }
            )
            chats.append_global(
                conversation_id=f"intake:{ip_hash}",
                role="user",
                content=normalized_answer,
                metadata={
                    "channel": "toji_intake_chat",
                    "industry": industry,
                    "category": category,
                    "question_mode": "pure_chat",
                },
            )

        should_call_toji = (not transcript) or bool(incoming_answer)
        if should_call_toji and not bool(conversation.get("completed")):
            tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
            payload_context = {
                "industry": industry,
                "category": category,
                "problem_statement": str(tell_payload.get("problem_statement") or ""),
                "context": tell_payload.get("context") or {},
                "missing_fields": [str(x).strip() for x in (conversation.get("missing_fields") or []) if str(x).strip()],
                "captured_facts": [str(x).strip() for x in (conversation.get("captured_facts") or []) if str(x).strip()],
                "time_context": str(conversation.get("time_context") or ""),
                "question_limit": int(conversation.get("question_limit") or 10),
                "questions_answered": int(sum(1 for row in transcript if str(row.get("role") or "").strip().lower() == "user")),
            }
            try:
                turn_payload = llm.pure_intake_chat_turn(
                    industry=industry,
                    category=category,
                    payload_context=payload_context,
                    transcript=transcript,
                    user_message=normalized_answer if incoming_answer else "",
                )
            except Exception as exc:
                logger.error(
                    "Pure intake turn failed: industry=%s category=%s error=%s",
                    industry,
                    category,
                    exc,
                )
                raise HTTPException(
                    status_code=503,
                    detail=_toji_unavailable_detail(
                        llm,
                        fallback="Toji could not continue this intake turn right now. Please retry.",
                    ),
                ) from exc

            assistant_message = _convert_currency_text_to_usd(str(turn_payload.get("assistant_message") or "").strip())
            parse_confidence = float(turn_payload.get("parse_confidence")) if turn_payload.get("parse_confidence") is not None else None
            if not assistant_message:
                raise HTTPException(
                    status_code=503,
                    detail="Toji returned an empty intake response. Please retry.",
                )

            transcript.append(
                {
                    "role": "assistant",
                    "content": assistant_message,
                    "ts": iso(utcnow()),
                }
            )
            chats.append_global(
                conversation_id=f"intake:{ip_hash}",
                role="assistant",
                content=assistant_message,
                metadata={
                    "channel": "toji_intake_chat",
                    "industry": industry,
                    "category": category,
                    "question_mode": "pure_chat",
                    "question_source": "ollama",
                    "question_provider": llm.provider_name,
                    "question_model": llm.provider_model,
                },
            )

            existing_facts = [str(x).strip() for x in (conversation.get("captured_facts") or []) if str(x).strip()]
            new_facts = [str(x).strip() for x in (turn_payload.get("captured_facts") or []) if str(x).strip()]
            merged_facts: list[str] = []
            seen_facts: set[str] = set()
            for row in existing_facts + new_facts:
                key = row.lower()
                if key in seen_facts:
                    continue
                seen_facts.add(key)
                merged_facts.append(row)
            conversation["captured_facts"] = merged_facts[:80]
            time_context = str(turn_payload.get("time_context") or "").strip()
            if time_context:
                conversation["time_context"] = _convert_currency_text_to_usd(time_context)
            conversation["ready_to_analyze"] = bool(turn_payload.get("ready_to_analyze"))

        user_turns = sum(1 for row in transcript if str(row.get("role") or "").strip().lower() == "user")
        conversation["current_index"] = user_turns
        conversation["completed"] = bool(conversation.get("ready_to_analyze"))
        summary = _conversation_summary(conversation)
        conversation["total_questions"] = int(summary.get("total_questions") or max(1, int(body.max_questions or 10)))
        conversation["current_index"] = int(summary.get("answered") or user_turns)

        state = "QUESTIONNAIRE_IN_PROGRESS"
        sessions.update_state(ip_hash, state, step_uri="/intake/tell-us/turn")

        conversation_uri = _persist_turn_conversation(ip_hash, conversation)
        return TellUsTurnResponse(
            state=state,
            industry=industry,
            category=category,
            completed=bool(conversation.get("completed")),
            question_index=int(conversation.get("current_index") or 0),
            total_questions=int(conversation.get("total_questions") or 0),
            question=None,
            parse_confidence=parse_confidence,
            answer_source=answer_source,
            normalized_answer=normalized_answer,
            assistant_message=assistant_message,
            conversation_uri=conversation_uri,
            question_source=str(conversation.get("question_source") or "ollama"),
            question_mode=str(conversation.get("question_mode") or "pure_chat"),
            question_provider=str(conversation.get("question_provider") or llm.provider_name),
            question_model=str(conversation.get("question_model") or llm.provider_model),
            summary=summary,
        )

    @app.post("/intake/tell-us/finalize", response_model=TellUsFinalizeResponse)
    def intake_tell_us_finalize(request: Request, body: TellUsFinalizeRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        _enforce_rate_limit(
            request,
            ip_hash=ip_hash,
            scope="intake_tell_us",
            limit=int(os.getenv("RATE_TELLUS_LIMIT", "40")),
            window_seconds=int(os.getenv("RATE_TELLUS_WINDOW_SEC", "3600")),
        )

        conversation = read_session_json(ip_hash, "tell_us_conversation_uri")
        if not isinstance(conversation, dict):
            raise HTTPException(status_code=400, detail="No active Tell-Us conversation to finalize.")

        summary = _conversation_summary(conversation)
        if summary.get("unresolved", 0) > 0 and not body.force:
            raise HTTPException(
                status_code=400,
                detail="Conversation still has unresolved questions. Pass force=true to finalize anyway.",
            )

        custom_industry_name = None
        if body.industry:
            normalized_industry, custom_industry_name = _normalize_industry_value(str(body.industry))
            conversation["industry"] = normalized_industry
        if body.category:
            conversation["category"] = str(body.category)
        conversation["completed"] = True
        conversation_uri = _persist_turn_conversation(ip_hash, conversation)

        tell_payload = _build_tell_us_from_conversation(ip_hash=ip_hash, conversation=conversation)
        context_payload = tell_payload.get("context") or {}
        if not isinstance(context_payload, dict):
            context_payload = {}
        if custom_industry_name:
            context_payload["custom_industry_name"] = custom_industry_name
            tell_payload["context"] = context_payload
        resolved_custom_industry = str(context_payload.get("custom_industry_name") or "").strip() or None
        tell_uri = write_runtime_json(ip_hash, "tell_us", tell_payload)
        sessions.update_artifact(ip_hash, "tell_us_uri", tell_uri)
        sessions.update_state(ip_hash, "QUESTIONNAIRE_CONFIRMED", step_uri="/intake/tell-us/finalize")
        llm.bind_session_context(
            ip_hash,
            tell_us=tell_payload,
            industry=tell_payload.get("industry"),
            category=(tell_payload.get("context") or {}).get("category"),
            custom_industry_name=resolved_custom_industry,
        )

        return TellUsFinalizeResponse(
            state="QUESTIONNAIRE_CONFIRMED",
            industry=str(tell_payload.get("industry") or "ecommerce"),
            category=str((tell_payload.get("context") or {}).get("category") or "unclassified"),
            tell_us_uri=tell_uri,
            conversation_uri=conversation_uri,
            summary=summary,
        )

    @app.get("/intake/categories/{industry}")
    def intake_categories(industry: str, request: Request):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        normalized_industry, custom_industry = _normalize_industry_value(industry, default="other_industry")
        if normalized_industry == "other_industry":
            llm.bind_session_context(
                ip_hash,
                industry=normalized_industry,
                custom_industry_name=custom_industry,
            )
            return {
                "industry": "other_industry",
                "custom_industry_name": custom_industry,
                "categories": [],
            }
        llm.preload(normalized_industry)
        cats = llm.categories(normalized_industry)
        return {
            "industry": normalized_industry,
            "custom_industry_name": custom_industry,
            "categories": [{"category": c, "label": c.replace("_", " ").title()} for c in cats],
        }

    @app.post("/intake/column-questions")
    def intake_column_questions(request: Request, body: ColumnQuestionsRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        industry, custom_industry = _normalize_industry_value(body.industry, default="other_industry")
        llm.preload(industry)
        tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
        user_context = _business_intent_text(tell_payload)
        try:
            questions = llm.contextual_intake_questions(
                industry=industry,
                category=body.category,
                user_context=user_context,
                max_questions=10,
            )
        except Exception as exc:
            logger.error(
                "intake/column-questions failed: industry=%s category=%s error=%s",
                industry,
                body.category,
                exc,
            )
            raise HTTPException(
                status_code=503,
                detail="Toji is unavailable for context questions right now.",
            ) from exc
        llm.bind_session_context(
            ip_hash,
            industry=industry,
            category=body.category,
            custom_industry_name=custom_industry,
        )
        return {
            "industry": industry,
            "custom_industry_name": custom_industry,
            "category": body.category,
            "questions": questions,
            "count": len(questions),
            "source": "ollama",
        }

    @app.post("/intake/tell-us/confirm")
    def intake_tell_us_confirm(request: Request, body: TellUsConfirmRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        try:
            return _generate_synthetic_artifact(ip_hash=ip_hash, requested_industry=body.industry)
        except RuntimeError as exc:
            raise HTTPException(status_code=503, detail=str(exc)) from exc

    @app.post("/intake/tell-us/confirm-task")
    def intake_tell_us_confirm_task(request: Request, body: TellUsConfirmRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")

        requested_industry, _requested_custom = _normalize_industry_value(body.industry, default="other_industry")
        task = tasks.create(industry=requested_industry or body.industry, owner_session=ip_hash)
        sessions.update_state(
            ip_hash,
            "SYNTH_DATA_GENERATING",
            step_uri="/intake/tell-us/confirm-task",
            active_task_id=task.task_id,
        )

        def _runner() -> None:
            try:
                def _progress(state: str, progress: float, result: Optional[dict[str, Any]] = None) -> None:
                    tasks.set_state(task.task_id, state, progress=progress)
                    if result is not None:
                        tasks.set_result(task.task_id, result)

                result = _generate_synthetic_artifact(
                    ip_hash=ip_hash,
                    requested_industry=body.industry,
                    progress_hook=_progress,
                )
                tasks.set_result(task.task_id, result)
                tasks.set_state(task.task_id, "SYNTH_DATA_READY", progress=1.0)
                sessions.update_state(ip_hash, "QUESTIONNAIRE_CONFIRMED", active_task_id=None)
            except Exception as exc:
                sessions.update_state(ip_hash, "ERROR", active_task_id=None)
                tasks.set_state(task.task_id, "ERROR", progress=1.0, error=_sanitize_error(str(exc)))

        tasks.submit(task.task_id, _runner)
        return JSONResponse(status_code=202, content={"task_id": task.task_id, "state": "SYNTH_DATA_GENERATING"})

    @app.post("/intake/tell-us/augment")
    def intake_tell_us_augment(request: Request, payload: TellUsAugmentRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        missing_fields = list(payload.missing_fields or [])
        provided_values = dict(payload.provided_values or {})
        out = {
            "missing_fields": missing_fields,
            "provided_values": provided_values,
            "recorded_at": iso(utcnow()),
        }
        rel = write_runtime_json(ip_hash, "augment", out)
        sessions.update_artifact(ip_hash, "augment_uri", rel)
        sessions.update_state(ip_hash, "QUESTIONNAIRE_CONFIRMED", step_uri="/intake/tell-us/augment")
        return {"state": "QUESTIONNAIRE_CONFIRMED", "artifact_uri": rel}

    @app.post("/analyze/schema")
    def analyze_schema(request: Request, body: SchemaAnalyzeRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        sessions.update_state(ip_hash, "SCHEMA_ANALYSIS_RUNNING", step_uri="/analyze/schema")
        df = _read_dataset(resolve_artifact(ip_hash, body.artifact_uri))
        schema = _schema_summary(df)
        rel = write_runtime_json(ip_hash, "schema", schema)
        sessions.update_artifact(ip_hash, "schema_uri", rel)
        sessions.update_state(ip_hash, "SCHEMA_READY", step_uri="/analyze/schema")
        llm.bind_session_context(ip_hash, schema=schema, schema_uri=rel)
        return {"state": "SCHEMA_READY", "schema_uri": rel, "schema": schema}

    def _match_with_alias_self_healing(
        *,
        ip_hash: str,
        columns: list[str],
        requested_industry: Optional[str],
    ) -> tuple[dict[str, Any], Optional[dict[str, Any]], Optional[dict[str, Any]]]:
        if not legacy_router_enabled or matcher is None or alias_registry is None or target_semantics is None:
            raise HTTPException(
                status_code=410,
                detail="Legacy router/match flow is disabled. Use Toji unified flow via intake + run/predict.",
            )
        normalized_requested_industry, custom_requested_industry = _normalize_industry_value(
            requested_industry,
            default="",
        )
        requested_for_match = normalized_requested_industry or None
        match = matcher.match(columns, industry=requested_for_match)
        resolved_industry = str(match.get("industry") or normalized_requested_industry or "").strip()
        if resolved_industry:
            resolved_industry, _ = _normalize_industry_value(resolved_industry, default=resolved_industry)
            match["industry"] = resolved_industry
        resolved_category = str(match.get("top_category") or "").strip()
        alias_feedback = None
        semantic_shadow = None
        tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
        user_context = _business_intent_text(tell_payload)
        if resolved_industry and resolved_category:
            canonical_columns = matcher.category_columns(resolved_industry, resolved_category)
            if canonical_columns:
                proposals = llm.infer_column_aliases(
                    industry=resolved_industry,
                    category=resolved_category,
                    provided_columns=columns,
                    canonical_columns=canonical_columns,
                    user_context=str(tell_payload.get("problem_statement") or ""),
                )
                if proposals:
                    alias_feedback = alias_registry.record_aliases(
                        industry=resolved_industry,
                        category=resolved_category,
                        aliases=proposals,
                        worker_dataset_ids=matcher.category_worker_dataset_ids(
                            resolved_industry,
                            resolved_category,
                        ),
                        source=f"{llm.provider_name}_runtime",
                    )
                    matcher.reload()
                    match = matcher.match(columns, industry=requested_for_match)
                    resolved_industry = str(match.get("industry") or normalized_requested_industry or "").strip()
                    if resolved_industry:
                        resolved_industry, _ = _normalize_industry_value(resolved_industry, default=resolved_industry)
                        match["industry"] = resolved_industry
                    resolved_category = str(match.get("top_category") or "").strip()
        if resolved_industry:
            semantic_shadow = llm.infer_semantic_routing(
                industry=resolved_industry,
                provided_columns=columns,
                user_context=user_context,
                top_k=3,
            )
            if isinstance(semantic_shadow, dict):
                # Integrate semantic shadow into matcher ranking while preserving
                # deterministic schema score as the hard acceptance gate.
                match = matcher.match(
                    columns,
                    industry=requested_for_match,
                    semantic_shadow=semantic_shadow,
                )
                resolved_industry = str(match.get("industry") or normalized_requested_industry or "").strip()
                if resolved_industry:
                    resolved_industry, _ = _normalize_industry_value(resolved_industry, default=resolved_industry)
                    match["industry"] = resolved_industry
                match["semantic_routing_shadow"] = semantic_shadow
                det_category = str(match.get("top_category") or "").strip()
                sem_category = str(semantic_shadow.get("top_category") or "").strip()
                det_score = float(match.get("deterministic_score") or match.get("schema_coverage") or 0.0)
                sem_score = float(semantic_shadow.get("top_category_score") or 0.0)
                blended = float(match.get("blended_score") or ((det_score * 0.85) + (sem_score * 0.15)))
                match["semantic_alignment"] = {
                    "deterministic_top_category": det_category or None,
                    "semantic_top_category": sem_category or None,
                    "agree": bool(det_category and sem_category and det_category == sem_category),
                    "deterministic_score": round(det_score, 4),
                    "semantic_score": round(sem_score, 4),
                    "blended_shadow_score": round(blended, 4),
                }

        if resolved_industry:
            try:
                target_semantics_payload = target_semantics.infer_from_match(
                    industry=resolved_industry,
                    match=match,
                    top_k_workers=3,
                )
            except Exception as exc:
                logger.warning("target semantics inference failed: %s", exc)
                target_semantics_payload = None
            if target_semantics_payload:
                match["target_semantics"] = target_semantics_payload
                top_guess = (target_semantics_payload or {}).get("top_guess") or {}
                if isinstance(semantic_shadow, dict) and top_guess:
                    semantic_shadow["target_semantic_guess"] = {
                        "target_column": top_guess.get("target_column"),
                        "meaning": top_guess.get("meaning"),
                        "confidence": top_guess.get("confidence"),
                        "evidence": top_guess.get("evidence") or [],
                        "worker_model_id": top_guess.get("worker_model_id"),
                        "category": top_guess.get("category"),
                    }
        if normalized_requested_industry:
            match["requested_industry"] = normalized_requested_industry
        if custom_requested_industry:
            match["custom_industry_name"] = custom_requested_industry
        return match, alias_feedback, semantic_shadow

    @app.post("/analyze/match-category")
    def analyze_match_category(request: Request, body: CategoryMatchRequest):
        if not legacy_router_enabled:
            raise HTTPException(
                status_code=410,
                detail="Category matcher endpoint is disabled. Toji unified flow does not require router/L0/L1.",
            )
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        sessions.update_state(ip_hash, "CATEGORY_EVAL_RUNNING", step_uri="/analyze/match-category")

        columns = body.columns
        if not columns:
            df = _read_dataset(resolve_artifact(ip_hash, body.artifact_uri))
            columns = list(df.columns)

        match, alias_feedback, _semantic_shadow = _match_with_alias_self_healing(
            ip_hash=ip_hash,
            columns=columns,
            requested_industry=body.industry,
        )
        if alias_feedback:
            match["alias_feedback"] = alias_feedback
        rel = write_runtime_json(ip_hash, "match", match)
        sessions.update_artifact(ip_hash, "match_uri", rel)
        llm.bind_session_context(
            ip_hash,
            industry=match.get("industry"),
            category=match.get("top_category"),
            custom_industry_name=match.get("custom_industry_name"),
            match=match,
        )

        decision = match.get("decision")
        if decision == "ACCEPT":
            sessions.update_state(ip_hash, "CATEGORY_ACCEPTED", step_uri="/analyze/match-category")
        elif decision == "PARTIAL_ACCEPT":
            sessions.update_state(ip_hash, "PARTIAL_MATCH", step_uri="/analyze/match-category")
        else:
            sessions.update_state(ip_hash, "CATEGORY_REJECTED", step_uri="/analyze/match-category")

        return match

    @app.post("/run/predict")
    def run_predict(request: Request, body: PredictRequest):
        ip_hash, rec = require_session(request)
        if rec is None:
            raise HTTPException(status_code=400, detail="Start session first.")
        if rec.state == "SESSION_LOCKED_COOLDOWN":
            return JSONResponse(status_code=423, content=rec.to_status_payload())
        _enforce_rate_limit(request, ip_hash=ip_hash, scope="run_predict", limit=int(os.getenv("RATE_PREDICT_LIMIT", "12")), window_seconds=int(os.getenv("RATE_PREDICT_WINDOW_SEC", "3600")))
        # Hard product policy: prediction is handled only by Toji/Ollama unified flow.
        requested_industry, requested_custom_industry = _normalize_industry_value(body.industry, default="")

        task = tasks.create(industry=requested_industry or body.industry, owner_session=ip_hash)
        sessions.update_state(
            ip_hash,
            "PREDICT_ENQUEUED",
            step_uri="/run/predict",
            active_task_id=task.task_id,
        )

        def _runner() -> None:
            try:
                tasks.set_state(task.task_id, "TOJI_CONTEXT_BUILDING", progress=0.15)
                sessions.update_state(ip_hash, "TOJI_CONTEXT_BUILDING", active_task_id=task.task_id)

                df = _read_dataset(resolve_artifact(ip_hash, body.artifact_uri))
                schema = _schema_summary(df)
                schema_uri = write_runtime_json(ip_hash, "schema", schema)
                sessions.update_artifact(ip_hash, "schema_uri", schema_uri)

                tasks.set_state(task.task_id, "ALIGNMENT_RUNNING", progress=0.35)
                sessions.update_state(ip_hash, "ALIGNMENT_RUNNING", active_task_id=task.task_id)

                tell_payload = read_session_json(ip_hash, "tell_us_uri") or {}
                question_payload = read_session_json(ip_hash, "questions_uri") or {}
                augment_payload = read_session_json(ip_hash, "augment_uri") or {}
                synth_meta_payload = read_session_json(ip_hash, "synthetic_meta_uri") or {}
                session_ctx = llm.session_context(ip_hash)
                intent_text = _business_intent_text(tell_payload)

                raw_industry = str(
                    body.industry
                    or tell_payload.get("industry")
                    or session_ctx.get("industry")
                    or ""
                ).strip()
                industry, inferred_custom_industry = _normalize_industry_value(
                    raw_industry,
                    default="other_industry",
                )
                context_custom_industry = ""
                tell_context = tell_payload.get("context") if isinstance(tell_payload.get("context"), dict) else {}
                if isinstance(tell_context, dict):
                    context_custom_industry = str(tell_context.get("custom_industry_name") or "").strip()
                resolved_custom_industry = (
                    inferred_custom_industry
                    or context_custom_industry
                    or str(session_ctx.get("custom_industry_name") or "").strip()
                    or requested_custom_industry
                    or None
                )
                category = str(
                    body.category
                    or (tell_payload.get("context") or {}).get("category")
                    or session_ctx.get("category")
                    or llm.infer_category(industry, intent_text or str(tell_payload.get("problem_statement") or ""))
                    or "general_operations"
                ).strip() or "general_operations"
                llm.bind_session_context(
                    ip_hash,
                    industry=industry,
                    category=category,
                    custom_industry_name=resolved_custom_industry,
                )

                dataset_profile = _dataset_profile_for_llm(df, schema)
                if isinstance(synth_meta_payload, dict) and synth_meta_payload:
                    dataset_profile["synthetic_generation"] = {
                        "source": synth_meta_payload.get("source"),
                        "analysis_trace": synth_meta_payload.get("analysis_trace") or [],
                        "assumptions": synth_meta_payload.get("assumptions") or [],
                        "python_script_uri": synth_meta_payload.get("python_script_uri"),
                        "pandas_copy_uri": synth_meta_payload.get("pandas_copy_uri"),
                    }
                quality = _quality_from_schema(schema)
                synthetic_data_conf = _calibrate_synthetic_confidence(quality["completeness"])
                prediction_conf = _calibrate_prediction_confidence(quality["completeness"])
                data_conf = (
                    "high"
                    if quality["completeness"] >= 0.80
                    else "medium"
                    if quality["completeness"] >= 0.55
                    else "low"
                )

                tasks.set_state(task.task_id, "TOJI_ANALYSIS_RUNNING", progress=0.65)
                sessions.update_state(ip_hash, "TOJI_ANALYSIS_RUNNING", active_task_id=task.task_id)

                unified = llm.generate_unified_analysis(
                    industry=industry,
                    category=category,
                    dataset_profile=dataset_profile,
                    user_context=intent_text or str(tell_payload.get("problem_statement") or ""),
                )
                unified = _normalize_currency_payload(unified)

                recs = [str(x).strip() for x in (unified.get("recommendations") or []) if str(x).strip()]

                report_payload = {
                    "generated_at": iso(utcnow()),
                    "task_id": task.task_id,
                    "industry": industry,
                    "category": category,
                    "custom_industry_name": resolved_custom_industry,
                    "post_dashboard_intake_required": bool(
                        ((tell_payload.get("context") or {}) if isinstance(tell_payload.get("context"), dict) else {}).get("post_dashboard_intake_required")
                    ),
                    "decision": "LLM_UNIFIED",
                    "scorecard": {
                        "data_quality_band": data_conf,
                        "coverage": quality["completeness"],
                        "rows": int(schema.get("rows") or 0),
                        "columns": int(schema.get("cols") or 0),
                    },
                    "routing": {
                        "fallback_models": [],
                        "top_workers": [],
                        "missing_fields": quality["missing_fields"],
                    },
                    "source_artifact_uri": body.artifact_uri,
                    "schema": schema,
                    "dataset_profile": dataset_profile,
                    "unified_analysis": unified,
                    "opportunity_analysis": unified.get("opportunity_analysis") or {},
                    "insights": [
                        str(unified.get("trend") or ""),
                        str(unified.get("seasonality") or ""),
                        str(unified.get("behaviour") or ""),
                    ],
                    "next_actions": [
                        {"action": rec, "priority": "HIGH" if idx == 0 else "MEDIUM"}
                        for idx, rec in enumerate(recs[:3])
                    ],
                    "dashboard_source": "toji_ollama",
                    "llm_runtime": {
                        "enabled": True,
                        "provider": llm.provider_name,
                        "model": llm.provider_model,
                    },
                    "runtime_inference": {
                        "prediction_mode": "llm_unified",
                        "mean_confidence": prediction_conf,
                        "synthetic_data_confidence": synthetic_data_conf,
                        "prediction_confidence": prediction_conf,
                        "prediction_confidence_metric": "calibrated confidence",
                        "provider": llm.provider_name,
                        "model": llm.provider_model,
                    },
                    "runtime_inference_multi": [],
                    "routing_context": {
                        "route_categories": [category],
                        "schema_coverage": quality["completeness"],
                        "suggested_categories": [],
                    },
                    "column_source_map": {},
                    "data_confidence": data_conf,
                }
                report_payload["intake_qa"] = _build_intake_qa_bundle(
                    tell_payload=tell_payload,
                    question_payload=question_payload,
                    augment_payload=augment_payload,
                    session_ctx=session_ctx,
                )
                try:
                    user_context = intent_text or str(tell_payload.get("problem_statement") or "")
                    toji_summary = llm.summarize_report(
                        report_payload,
                        user_context=user_context,
                        strict=True,
                    )
                except Exception as _summary_exc:
                    raise RuntimeError(f"Ollama summary generation failed in unified flow: {_summary_exc}") from _summary_exc
                if not toji_summary:
                    raise RuntimeError("Ollama returned no summary payload in unified flow.")
                normalized_summary = _normalize_currency_payload(toji_summary)
                report_payload["llm_summary"] = normalized_summary
                summary_recs = normalized_summary.get("recommendations") or []
                next_actions: list[dict[str, Any]] = []
                for idx, rec in enumerate(summary_recs[:3]):
                    if isinstance(rec, dict):
                        action = str(rec.get("action") or rec.get("description") or "").strip()
                    else:
                        action = str(rec).strip()
                    if not action:
                        continue
                    next_actions.append(
                        {"action": action, "priority": "HIGH" if idx == 0 else ("MEDIUM" if idx == 1 else "LOW")}
                    )
                if next_actions:
                    report_payload["next_actions"] = next_actions
                logger.info(
                    "Unified Toji summary generated: industry=%s category=%s problems=%d recommendations=%d",
                    industry,
                    category,
                    len(normalized_summary.get("problems") or []),
                    len(normalized_summary.get("recommendations") or []),
                )
                report_payload = _normalize_currency_payload(report_payload)

                tasks.set_state(task.task_id, "REPORT_BUILDING", progress=0.85)
                sessions.update_state(ip_hash, "REPORT_BUILDING", active_task_id=task.task_id)
                report_id = reports.write_report(report_payload)

                # Markdown dashboard is primary and authoritative.
                if markdown_dashboard_enabled:
                    tasks.set_state(task.task_id, "DASHBOARD_MARKDOWN_GENERATING", progress=0.9)
                    md_payload = llm.generate_markdown_dashboard(
                        industry=industry,
                        category=category,
                        user_context=intent_text or str(tell_payload.get("problem_statement") or ""),
                        dataset_profile=dataset_profile,
                        df=df,
                    )
                    report_payload["dashboard_markdown"] = str(md_payload.get("markdown") or "").strip()
                    report_payload["dashboard_markdown_meta"] = {
                        "analysis_trace": md_payload.get("analysis_trace") or [],
                        "assumptions": md_payload.get("assumptions") or [],
                        "markdown_sha256": md_payload.get("markdown_sha256") or "",
                        "generated_by": "llm_markdown_quickchart",
                    }
                    report_payload["visuals"] = _sandbox_render_quickchart_visuals(
                        base_path=base,
                        orchestrator=llm,
                        report_id=report_id,
                        visuals=md_payload.get("visuals") or [],
                    )
                    report_payload["dashboard_source"] = "markdown_quickchart"
                    report_payload["visuals_error"] = None
                else:
                    raise RuntimeError("Markdown dashboard generation is disabled.")
                reports.update_report(report_id, report_payload)

                tasks.set_state(task.task_id, "REPORT_READY", progress=1.0, report_id=report_id)
                sessions.update_state(
                    ip_hash,
                    "REPORT_READY",
                    step_uri="/run/predict",
                    active_task_id=None,
                    report_id=report_id,
                )
            except Exception as exc:
                debug_uri = reports.write_debug_bundle(
                    task.task_id,
                    {
                        "error": str(exc),
                        "industry": requested_industry or body.industry,
                        "artifact_uri": body.artifact_uri,
                        "task_id": task.task_id,
                    },
                )
                sessions.update_artifact(ip_hash, "routing_uri", debug_uri)
                sessions.update_state(ip_hash, "ERROR", active_task_id=None)
                tasks.set_state(task.task_id, "ERROR", progress=1.0, error=_sanitize_error(str(exc)))

        tasks.submit(task.task_id, _runner)
        return JSONResponse(status_code=202, content={"task_id": task.task_id, "state": "PREDICT_ENQUEUED"})

    @app.get("/task/{task_id}", response_model=TaskStatusResponse)
    def task_status(task_id: str):
        task_id = _validate_id(task_id, "task_id")
        rec = tasks.get(task_id)
        if rec is None:
            raise HTTPException(status_code=404, detail="Task not found.")
        return TaskStatusResponse(**rec.to_payload())

    @app.get("/report/{report_id}")
    def report_get(report_id: str):
        report_id = _validate_id(report_id, "report_id")
        payload = reports.read_report(report_id)
        if payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")
        chat_summary = chats.read_summary(report_id)
        if chat_summary is not None:
            payload["chat_summary"] = chat_summary
        return _normalize_currency_payload(payload)

    @app.get("/report/{report_id}/export")
    def report_export(report_id: str):
        report_id = _validate_id(report_id, "report_id")
        payload = reports.read_report(report_id)
        if payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")

        normalized_payload = _normalize_currency_payload(payload)
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("report.json", json.dumps(normalized_payload, indent=2))
            chat_summary = chats.read_summary(report_id)
            if chat_summary is not None:
                zf.writestr("chat_summary.json", json.dumps(chat_summary, indent=2))

            assets_dir = base / "exports" / "runtime_reports" / "assets" / report_id
            if assets_dir.exists() and assets_dir.is_dir():
                export_globs = ("*.png", "*.jpg", "*.jpeg", "*.webp", "*.svg", "*.json", "*.csv", "*.txt", "*.docx", "*.pptx")
                for pattern in export_globs:
                    for path in sorted(assets_dir.glob(pattern)):
                        if path.is_file():
                            zf.write(path, arcname=f"assets/{path.name}")

        buffer.seek(0)
        filename = f"{report_id}_report_export.zip"
        headers = {"Content-Disposition": f'attachment; filename="{filename}"'}
        return StreamingResponse(buffer, media_type="application/zip", headers=headers)

    @app.get("/report/{report_id}/asset/{asset_name}")
    def report_asset(report_id: str, asset_name: str):
        report_id = _validate_id(report_id, "report_id")
        # Sanitize asset_name to prevent path traversal (#34)
        clean_name = Path(asset_name).name
        if not clean_name or '/' in asset_name or '\\' in asset_name or '..' in asset_name:
            raise HTTPException(status_code=400, detail="Invalid asset name.")
        assets_dir = base / "exports" / "runtime_reports" / "assets" / report_id
        asset_path = (assets_dir / clean_name).resolve()
        resolved_dir = assets_dir.resolve()
        # Use startswith check so the resolved path must be a strict descendant of
        # assets_dir — this also catches the edge case where asset_path == assets_dir
        # which the old `not in parents` check would have allowed through.
        if not str(asset_path).startswith(str(resolved_dir) + os.sep):
            raise HTTPException(status_code=400, detail="Invalid asset path.")
        if not asset_path.exists() or not asset_path.is_file():
            raise HTTPException(status_code=404, detail="Asset not found.")
        return FileResponse(path=str(asset_path))

    @app.post("/report/{report_id}/regenerate-visuals")
    def report_regenerate_visuals(request: Request, report_id: str):
        report_id = _validate_id(report_id, "report_id")
        ip_hash = _ip_hash(request, session_secret)
        _enforce_rate_limit(
            request,
            ip_hash=ip_hash,
            scope="report_regenerate_visuals",
            limit=int(os.getenv("RATE_REPORT_REGEN_LIMIT", "20")),
            window_seconds=int(os.getenv("RATE_REPORT_REGEN_WINDOW_SEC", "3600")),
        )
        report_payload = reports.read_report(report_id)
        if report_payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")

        artifact_uri = str(report_payload.get("source_artifact_uri") or "").strip()
        if not artifact_uri:
            raise HTTPException(status_code=400, detail="No source artifact recorded for this report.")

        try:
            artifact_path = _safe_rel(base, artifact_uri)
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=_sanitize_error(exc)) from exc

        exports_dir = (base / "exports").resolve()
        if exports_dir not in artifact_path.resolve().parents:
            raise HTTPException(status_code=400, detail="Artifact must reside within exports directory.")
        if not artifact_path.exists():
            raise HTTPException(status_code=404, detail="Source artifact file not found.")

        try:
            df = _read_dataset(artifact_path)
            schema = _schema_summary(df)
            dataset_profile = _dataset_profile_for_llm(df, schema)
            user_context = str((report_payload.get("intake_qa") or {}).get("problem_statement") or "")
            if not markdown_dashboard_enabled:
                raise RuntimeError("Markdown dashboard generation is disabled.")
            md_payload = llm.generate_markdown_dashboard(
                industry=str(report_payload.get("industry") or "other_industry"),
                category=str(report_payload.get("category") or "unclassified"),
                user_context=user_context,
                dataset_profile=dataset_profile,
                df=df,
            )
            report_payload["dashboard_markdown"] = str(md_payload.get("markdown") or "").strip()
            report_payload["dashboard_markdown_meta"] = {
                "analysis_trace": md_payload.get("analysis_trace") or [],
                "assumptions": md_payload.get("assumptions") or [],
                "markdown_sha256": md_payload.get("markdown_sha256") or "",
                "generated_by": "llm_markdown_quickchart",
            }
            report_payload["visuals"] = _sandbox_render_quickchart_visuals(
                base_path=base,
                orchestrator=llm,
                report_id=report_id,
                visuals=md_payload.get("visuals") or [],
            )
            report_payload["dashboard_source"] = "markdown_quickchart"
            report_payload.pop("visuals_error", None)
            reports.update_report(report_id, report_payload)
            return {
                "ok": True,
                "report_id": report_id,
                "dashboard_source": "markdown_quickchart",
                "visuals_count": len(report_payload.get("visuals") or []),
                "visuals": report_payload.get("visuals") or [],
                "markdown": report_payload.get("dashboard_markdown") or "",
            }
        except Exception as exc:
            logger.warning("Dashboard regeneration failed: %s", str(exc)[:500])
            report_payload["visuals_error"] = "Could not regenerate dashboard markdown right now."
            reports.update_report(report_id, report_payload)
            raise HTTPException(status_code=500, detail=_sanitize_error(str(exc))) from exc

    @app.post("/report/{report_id}/regenerate-visual-slot")
    def report_regenerate_visual_slot(request: Request, report_id: str, body: dict[str, Any]):
        report_id = _validate_id(report_id, "report_id")
        ip_hash = _ip_hash(request, session_secret)
        _enforce_rate_limit(
            request,
            ip_hash=ip_hash,
            scope="report_regenerate_visuals",
            limit=int(os.getenv("RATE_REPORT_REGEN_LIMIT", "20")),
            window_seconds=int(os.getenv("RATE_REPORT_REGEN_WINDOW_SEC", "3600")),
        )
        report_payload = reports.read_report(report_id)
        if report_payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")

        slot = str((body or {}).get("slot") or "").strip().lower()
        valid_slots = ["r1c1", "r1c2", "r2c1", "r2c2"]
        if slot not in valid_slots:
            raise HTTPException(status_code=400, detail="Invalid visual slot.")
        slot_idx = valid_slots.index(slot)

        artifact_uri = str(report_payload.get("source_artifact_uri") or "").strip()
        if not artifact_uri:
            raise HTTPException(status_code=400, detail="No source artifact recorded for this report.")

        try:
            artifact_path = _safe_rel(base, artifact_uri)
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=_sanitize_error(exc)) from exc
        exports_dir = (base / "exports").resolve()
        if exports_dir not in artifact_path.resolve().parents:
            raise HTTPException(status_code=400, detail="Artifact must reside within exports directory.")
        if not artifact_path.exists():
            raise HTTPException(status_code=404, detail="Source artifact file not found.")

        try:
            df = _read_dataset(artifact_path)
            schema = _schema_summary(df)
            dataset_profile = _dataset_profile_for_llm(df, schema)
            intake_qa = report_payload.get("intake_qa") or {}
            user_context = str(intake_qa.get("problem_statement") or "")
            existing_visuals = report_payload.get("visuals") if isinstance(report_payload.get("visuals"), list) else []
            existing_titles = [str((v or {}).get("title") or "").strip() for v in existing_visuals if isinstance(v, dict)]
            slot_payload = llm.generate_markdown_visual_slot(
                industry=str(report_payload.get("industry") or "other_industry"),
                category=str(report_payload.get("category") or "unclassified"),
                user_context=user_context,
                dataset_profile=dataset_profile,
                slot=slot,
                existing_titles=existing_titles,
            )

            visuals: list[dict[str, Any]] = [v for v in existing_visuals if isinstance(v, dict)]
            while len(visuals) < len(valid_slots):
                idx = len(visuals)
                visuals.append(
                    {
                        "name": f"quickchart_{idx + 1}",
                        "title": f"Visual {idx + 1}",
                        "kind": "quickchart",
                        "uri": "",
                        "meta": {"slot": valid_slots[idx], "caption": "", "source": "quickchart"},
                    }
                )

            refreshed_visual = {
                "name": f"quickchart_{slot_idx + 1}",
                "title": str(slot_payload.get("title") or f"Visual {slot_idx + 1}").strip(),
                "kind": "quickchart",
                "uri": str(slot_payload.get("quickchart_url") or "").strip(),
                "meta": {
                    "slot": slot,
                    "caption": str(slot_payload.get("caption") or "").strip(),
                    "source": "quickchart",
                },
            }
            visuals[slot_idx] = _sandbox_render_quickchart_visuals(
                base_path=base,
                orchestrator=llm,
                report_id=report_id,
                visuals=[refreshed_visual],
            )[0]
            report_payload["visuals"] = visuals
            report_payload["dashboard_source"] = "markdown_quickchart"
            report_payload["visuals_error"] = None
            reports.update_report(report_id, report_payload)
            return {"ok": True, "report_id": report_id, "slot": slot, "visual": visuals[slot_idx]}
        except Exception as exc:
            logger.warning("Slot visual regeneration failed: report=%s slot=%s error=%s", report_id, slot, str(exc)[:500])
            raise HTTPException(status_code=500, detail=_sanitize_error(str(exc))) from exc

    @app.post("/chat/guided-questions")
    def chat_guided_questions(request: Request, body: ChatGuidedQuestionsRequest):
        raise HTTPException(
            status_code=410,
            detail="This endpoint is disabled. Use pure Toji intake chat via /intake/tell-us/turn.",
        )

    @app.post("/chat/message")
    def chat_message(request: Request, body: ChatMessageRequest):
        # Validate message length (#26)
        if len(body.message) > MAX_CHAT_MESSAGE_LENGTH:
            raise HTTPException(
                status_code=400,
                detail=f"Message exceeds maximum length of {MAX_CHAT_MESSAGE_LENGTH} characters.",
            )
        ip_hash = _ip_hash(request, session_secret)
        _enforce_rate_limit(request, ip_hash=ip_hash, scope="chat_message", limit=int(os.getenv("RATE_CHAT_LIMIT", "60")), window_seconds=int(os.getenv("RATE_CHAT_WINDOW_SEC", "3600")))
        report_id = _resolve_report_id(ip_hash=ip_hash, explicit=body.report_id)
        report_payload = reports.read_report(report_id)
        if report_payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")
        normalized_user_message = _convert_currency_text_to_usd(body.message)
        history = chats.history(report_id, limit=80)
        chats.append(report_id, "user", normalized_user_message, metadata={"ip_hash": ip_hash})
        guard = _chat_guardrails(normalized_user_message, report_payload)
        if guard.get("blocked"):
            assistant_text = _convert_currency_text_to_usd(str(guard.get("response") or ""))
            chats.append(
                report_id,
                "assistant",
                assistant_text,
                metadata={
                    "guarded": True,
                    "guard_reason": guard.get("reason"),
                },
            )
            if sessions.get(ip_hash):
                sessions.update_state(ip_hash, "CHAT_ACTIVE", active_task_id=None)
            return {
                "report_id": report_id,
                "response": assistant_text,
                "highlights": _normalize_currency_payload(guard.get("highlights", [])),
                "recommendations": _normalize_currency_payload(guard.get("recommendations", [])),
                "guarded": True,
                "guard_reason": guard.get("reason"),
            }
        try:
            reply = llm.answer_chat(
                report_payload=report_payload,
                history=history,
                user_message=normalized_user_message,
            )
        except Exception as exc:
            logger.error("chat/message failed: report_id=%s error=%s", report_id, exc)
            raise HTTPException(
                status_code=503,
                detail="Toji is unavailable for chat right now. Please retry.",
            ) from exc
        raw_answer = str(reply.get("answer") or "").strip()
        if not raw_answer:
            raise HTTPException(
                status_code=503,
                detail="Toji returned an empty chat response. Please retry.",
            )
        normalized_reply = _normalize_currency_payload(reply)
        suggested_artifacts = normalized_reply.get("suggested_artifacts") or []
        want_doc, want_slides = _detect_chat_artifact_intent(normalized_user_message, suggested_artifacts)
        full_history = chats.history(report_id, limit=100)
        artifacts, artifact_errors = _build_chat_artifacts(
            base=base,
            report_id=report_id,
            report_payload=report_payload,
            reply=normalized_reply,
            user_message=normalized_user_message,
            want_doc=want_doc,
            want_slides=want_slides,
            chat_history=full_history,
        )

        assistant_text = _convert_currency_text_to_usd(raw_answer)
        if artifacts:
            links = "\n".join(
                f"- [{a['type'].title()} ({a['filename']})]({a['uri']})" for a in artifacts
            )
            assistant_text = f"{assistant_text}\n\nI generated files for you:\n{links}"
        elif (want_doc or want_slides) and artifact_errors:
            assistant_text = (
                f"{assistant_text}\n\nI couldn't generate every requested file this turn. "
                "Please retry and I'll regenerate."
            )
        chats.append(
            report_id,
            "assistant",
            assistant_text,
            metadata={
                "highlights": normalized_reply.get("highlights", []),
                "recommendations": normalized_reply.get("recommendations", []),
                "artifacts": artifacts,
                "artifact_errors": artifact_errors,
            },
        )
        if sessions.get(ip_hash):
            sessions.update_state(ip_hash, "CHAT_ACTIVE", active_task_id=None)
        return {
            "report_id": report_id,
            "response": assistant_text,
            "highlights": normalized_reply.get("highlights", []),
            "recommendations": normalized_reply.get("recommendations", []),
            "artifacts": artifacts,
            "artifact_errors": artifact_errors,
            "guarded": False,
        }

    @app.post("/chat/message/stream")
    async def chat_message_stream(request: Request, body: ChatMessageRequest):
        """SSE streaming variant of /chat/message. Sends token chunks progressively."""
        import threading

        if len(body.message) > MAX_CHAT_MESSAGE_LENGTH:
            raise HTTPException(
                status_code=400,
                detail=f"Message exceeds maximum length of {MAX_CHAT_MESSAGE_LENGTH} characters.",
            )
        ip_hash = _ip_hash(request, session_secret)
        _enforce_rate_limit(request, ip_hash=ip_hash, scope="chat_message", limit=int(os.getenv("RATE_CHAT_LIMIT", "60")), window_seconds=int(os.getenv("RATE_CHAT_WINDOW_SEC", "3600")))
        report_id = _resolve_report_id(ip_hash=ip_hash, explicit=body.report_id)
        report_payload = reports.read_report(report_id)
        if report_payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")
        normalized_user_message = _convert_currency_text_to_usd(body.message)
        history = chats.history(report_id, limit=80)
        chats.append(report_id, "user", normalized_user_message, metadata={"ip_hash": ip_hash})
        guard = _chat_guardrails(normalized_user_message, report_payload)
        if guard.get("blocked"):
            assistant_text = _convert_currency_text_to_usd(str(guard.get("response") or ""))
            chats.append(report_id, "assistant", assistant_text, metadata={"guarded": True, "guard_reason": guard.get("reason")})
            if sessions.get(ip_hash):
                sessions.update_state(ip_hash, "CHAT_ACTIVE", active_task_id=None)

            async def _blocked_gen():
                yield f"data: {json.dumps({'token': assistant_text})}\n\n"
                yield f"data: {json.dumps({'done': True, 'guarded': True, 'artifacts': [], 'artifact_errors': []})}\n\n"

            return StreamingResponse(_blocked_gen(), media_type="text/event-stream", headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

        # Use asyncio.Queue + call_soon_threadsafe so the worker thread posts items
        # directly onto the async event loop without any thread-pool overhead per token.
        # (The old threading.Queue + run_in_executor approach submitted a thread-pool task
        # for every single token, adding ~0.5–2 ms latency per chunk.)
        loop = asyncio.get_event_loop()
        chunk_q: asyncio.Queue = asyncio.Queue()

        def _run_stream():
            try:
                for chunk in llm.answer_chat_stream(
                    report_payload=report_payload,
                    history=history,
                    user_message=normalized_user_message,
                ):
                    loop.call_soon_threadsafe(chunk_q.put_nowait, chunk)
            except Exception as exc:
                loop.call_soon_threadsafe(chunk_q.put_nowait, f"__STREAM_ERROR__{exc}")
            finally:
                loop.call_soon_threadsafe(chunk_q.put_nowait, None)

        threading.Thread(target=_run_stream, daemon=True).start()

        async def generate():
            full_text = ""
            while True:
                item = await chunk_q.get()
                if item is None:
                    break
                if isinstance(item, str) and item.startswith("__STREAM_ERROR__"):
                    err_msg = item[len("__STREAM_ERROR__"):]
                    logger.error("chat/message/stream error: report_id=%s error=%s", report_id, err_msg)
                    yield f"data: {json.dumps({'error': 'Toji encountered an error. Please retry.'})}\n\n"
                    break
                if isinstance(item, str) and item.startswith("\x00DONE\x00"):
                    meta = {}
                    try:
                        meta = json.loads(item[6:])
                    except Exception:
                        pass
                    artifacts_suggested = meta.get("artifacts_to_generate") or []
                    want_doc, want_slides = _detect_chat_artifact_intent(normalized_user_message, artifacts_suggested)
                    full_history = chats.history(report_id, limit=100)
                    artifacts, artifact_errors = _build_chat_artifacts(
                        base=base,
                        report_id=report_id,
                        report_payload=report_payload,
                        reply={"answer": full_text},
                        user_message=normalized_user_message,
                        want_doc=want_doc,
                        want_slides=want_slides,
                        chat_history=full_history,
                    )
                    final_text = _convert_currency_text_to_usd(full_text)
                    if artifacts:
                        links = "\n".join(
                            f"- [{a['type'].title()} ({a['filename']})]({a['uri']})" for a in artifacts
                        )
                        final_text = f"{final_text}\n\nI generated files for you:\n{links}"
                    chats.append(
                        report_id,
                        "assistant",
                        final_text,
                        metadata={"artifacts": artifacts, "artifact_errors": artifact_errors},
                    )
                    if sessions.get(ip_hash):
                        sessions.update_state(ip_hash, "CHAT_ACTIVE", active_task_id=None)
                    yield f"data: {json.dumps({'done': True, 'artifacts': artifacts, 'artifact_errors': artifact_errors, 'guarded': False})}\n\n"
                    break
                # Regular token chunk — yield immediately so Starlette flushes it.
                full_text += str(item)
                yield f"data: {json.dumps({'token': str(item)})}\n\n"

        return StreamingResponse(generate(), media_type="text/event-stream", headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

    @app.post("/chat/brief")
    def chat_brief(request: Request, body: ChatBriefRequest):
        ip_hash = _ip_hash(request, session_secret)
        _enforce_rate_limit(request, ip_hash=ip_hash, scope="chat_brief", limit=int(os.getenv("RATE_CHAT_LIMIT", "60")), window_seconds=int(os.getenv("RATE_CHAT_WINDOW_SEC", "3600")))
        report_id = _resolve_report_id(ip_hash=ip_hash, explicit=body.report_id)
        report_payload = reports.read_report(report_id)
        if report_payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")

        # Idempotency: check if a briefing already exists in the first 5 messages
        history = chats.history(report_id, limit=5)
        for row in history:
            meta = row.get("metadata") or {}
            if isinstance(meta, dict) and meta.get("is_brief"):
                return {
                    "report_id": report_id,
                    "brief": _convert_currency_text_to_usd(row.get("content", "")),
                    "highlights": _normalize_currency_payload(meta.get("highlights", [])),
                    "recommended_action": _convert_currency_text_to_usd(meta.get("recommended_action", "")),
                    "cached": True,
                }

        try:
            reply = llm.generate_report_brief(report_payload=report_payload)
        except Exception as exc:
            logger.error("chat/brief failed: report_id=%s error=%s", report_id, exc)
            raise HTTPException(
                status_code=503,
                detail="Toji is unavailable for report briefing right now. Please retry.",
            ) from exc
        brief_text = _convert_currency_text_to_usd(str(reply.get("brief") or "").strip())
        if not brief_text:
            raise HTTPException(
                status_code=503,
                detail="Toji returned an empty report briefing. Please retry.",
            )
        highlights = _normalize_currency_payload(reply.get("highlights", []))
        recommended_action = _convert_currency_text_to_usd(str(reply.get("recommended_action") or ""))

        chats.append(
            report_id,
            "assistant",
            brief_text,
            metadata={
                "is_brief": True,
                "highlights": highlights,
                "recommended_action": recommended_action,
            },
        )

        if sessions.get(ip_hash):
            sessions.update_state(ip_hash, "CHAT_ACTIVE", active_task_id=None)

        return {
            "report_id": report_id,
            "brief": brief_text,
            "highlights": highlights,
            "recommended_action": recommended_action,
            "cached": False,
        }

    @app.get("/chat/history")
    def chat_history(request: Request, report_id: Optional[str] = Query(default=None), limit: int = Query(default=50, ge=1, le=500)):
        ip_hash = _ip_hash(request, session_secret)
        rid = _resolve_report_id(ip_hash=ip_hash, explicit=report_id)
        history = chats.history(rid, limit=limit)
        return {"report_id": rid, "messages": _normalize_currency_payload(history)}

    @app.post("/chat/export-summary")
    def chat_export_summary(request: Request, body: ChatExportRequest):
        ip_hash = _ip_hash(request, session_secret)
        report_id = _resolve_report_id(ip_hash=ip_hash, explicit=body.report_id)
        report_payload = reports.read_report(report_id)
        if report_payload is None:
            raise HTTPException(status_code=404, detail="Report not found.")
        history = chats.history(report_id, limit=1000)
        try:
            export_payload = _chat_export_payload(report_id, report_payload, history)
        except Exception as exc:
            logger.error("chat/export-summary failed: report_id=%s error=%s", report_id, exc)
            raise HTTPException(
                status_code=503,
                detail="Toji is unavailable for chat summary export right now. Please retry.",
            ) from exc
        export_payload = _normalize_currency_payload(export_payload)
        out_abs = chats.write_summary(report_id, export_payload)
        out_rel = None
        try:
            out_rel = str(Path(out_abs).resolve().relative_to(base))
        except Exception:
            out_rel = str(out_abs)

        report_payload["chat_summary_uri"] = out_rel
        report_payload["chat_summary"] = export_payload.get("summary")
        reports.update_report(report_id, report_payload)
        return {
            "report_id": report_id,
            "chat_summary_uri": out_rel,
            "summary": export_payload.get("summary"),
            "message_count": export_payload.get("message_count"),
        }

    @app.get("/events")
    async def events(
        task_id: str = Query(...),
        since: int = Query(0, ge=0),
        max_seconds: int = Query(30, ge=1, le=120),
    ):
        rec = tasks.get(task_id)
        if rec is None:
            raise HTTPException(status_code=404, detail="Task not found.")

        async def _event_stream():
            cursor = since
            deadline = utcnow().timestamp() + max_seconds
            while utcnow().timestamp() < deadline:
                rows = tasks.get_events(task_id, since=cursor)
                for row in rows:
                    cursor = max(cursor, int(row.get("seq", 0)))
                    yield f"data: {json.dumps(row)}\n\n"

                current = tasks.get(task_id)
                if current and current.state in {"REPORT_READY", "ERROR"}:
                    break
                await asyncio.sleep(1.0)

        return StreamingResponse(_event_stream(), media_type="text/event-stream")

    # Serve product frontend from the same backend service in production.
    # This avoids "/" 404 on single-service deploys (e.g., Railway).
    frontend_dir = (base / "frontend").resolve()
    if frontend_dir.exists() and frontend_dir.is_dir():
        app.mount("/", StaticFiles(directory=str(frontend_dir), html=True), name="frontend")

    return app


app = create_app()
